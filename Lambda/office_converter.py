"""
Office document converter - converts PPTX, DOCX, XLSX to text and images
Zero cost - uses only free Python libraries
"""
import io
import logging
from PIL import Image
from pptx import Presentation
from docx import Document as DocxDocument
from docx.shared import Pt
from openpyxl import load_workbook

logger = logging.getLogger(__name__)

def extract_pptx(file_path):
    """Extract text and images from PowerPoint"""
    text_content = []
    images = []
    
    try:
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract images
                if shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        if len(image_bytes) > 10000:  # Min 10KB
                            img_pil = Image.open(io.BytesIO(image_bytes))
                            img_byte_arr = io.BytesIO()
                            img_pil.convert('RGB').save(img_byte_arr, format='JPEG', quality=90)
                            images.append({
                                'page': slide_num,
                                'index': len(images),
                                'data': img_byte_arr.getvalue(),
                                'ext': 'jpg'
                            })
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide {slide_num}: {e}")
            
            if slide_text:
                text_content.append(f"[SLIDE {slide_num}]:\n" + "\n".join(slide_text))
        
        full_text = "\n\n".join(text_content)
        logger.info(f"PPTX: Extracted {len(prs.slides)} slides, {len(images)} images, {len(full_text)} chars")
        return full_text, images
        
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return "", []

def extract_docx(file_path):
    """Extract text and images from Word document in document order"""
    text_content = []
    images = []
    
    try:
        doc = DocxDocument(file_path)
        seen_rIds = set()  # Track relationship IDs instead of blobs to avoid missing images
        
        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text.strip())
        
        # Extract images in document order by walking through all elements
        import re
        
        # Walk through paragraphs in order
        for para in doc.paragraphs:
            for run in para.runs:
                if 'graphic' in run._element.xml:
                    rId_match = re.search(r'r:embed="(rId\d+)"', run._element.xml)
                    if rId_match:
                        rId = rId_match.group(1)
                        if rId not in seen_rIds:
                            seen_rIds.add(rId)
                            try:
                                rel = run.part.rels[rId]
                                if "image" in rel.target_ref:
                                    image_bytes = rel.target_part.blob
                                    if len(image_bytes) > 10000:
                                        img_pil = Image.open(io.BytesIO(image_bytes))
                                        img_byte_arr = io.BytesIO()
                                        img_pil.convert('RGB').save(img_byte_arr, format='JPEG', quality=90)
                                        images.append({
                                            'page': None,
                                            'index': len(images),
                                            'data': img_byte_arr.getvalue(),
                                            'ext': 'jpg'
                                        })
                            except Exception as e:
                                logger.warning(f"Failed to extract image from paragraph: {e}")
        
        # Walk through tables in order
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if 'graphic' in run._element.xml:
                                rId_match = re.search(r'r:embed="(rId\d+)"', run._element.xml)
                                if rId_match:
                                    rId = rId_match.group(1)
                                    if rId not in seen_rIds:
                                        seen_rIds.add(rId)
                                        try:
                                            rel = run.part.rels[rId]
                                            if "image" in rel.target_ref:
                                                image_bytes = rel.target_part.blob
                                                if len(image_bytes) > 10000:
                                                    img_pil = Image.open(io.BytesIO(image_bytes))
                                                    img_byte_arr = io.BytesIO()
                                                    img_pil.convert('RGB').save(img_byte_arr, format='JPEG', quality=90)
                                                    images.append({
                                                        'page': None,
                                                        'index': len(images),
                                                        'data': img_byte_arr.getvalue(),
                                                        'ext': 'jpg'
                                                    })
                                        except Exception as e:
                                            logger.warning(f"Failed to extract image from table: {e}")
        
        # Check headers/footers last
        for section in doc.sections:
            for header in [section.header, section.first_page_header, section.even_page_header]:
                try:
                    for para in header.paragraphs:
                        for run in para.runs:
                            if 'graphic' in run._element.xml:
                                rId_match = re.search(r'r:embed="(rId\d+)"', run._element.xml)
                                if rId_match:
                                    rId = rId_match.group(1)
                                    if rId not in seen_rIds:
                                        seen_rIds.add(rId)
                                        try:
                                            rel = run.part.rels[rId]
                                            if "image" in rel.target_ref:
                                                image_bytes = rel.target_part.blob
                                                if len(image_bytes) > 10000:
                                                    img_pil = Image.open(io.BytesIO(image_bytes))
                                                    img_byte_arr = io.BytesIO()
                                                    img_pil.convert('RGB').save(img_byte_arr, format='JPEG', quality=90)
                                                    images.append({
                                                        'page': None,
                                                        'index': len(images),
                                                        'data': img_byte_arr.getvalue(),
                                                        'ext': 'jpg'
                                                    })
                                        except Exception as e:
                                            logger.warning(f"Failed to extract image from header: {e}")
                except:
                    pass
            
            for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
                try:
                    for para in footer.paragraphs:
                        for run in para.runs:
                            if 'graphic' in run._element.xml:
                                rId_match = re.search(r'r:embed="(rId\d+)"', run._element.xml)
                                if rId_match:
                                    rId = rId_match.group(1)
                                    if rId not in seen_rIds:
                                        seen_rIds.add(rId)
                                        try:
                                            rel = run.part.rels[rId]
                                            if "image" in rel.target_ref:
                                                image_bytes = rel.target_part.blob
                                                if len(image_bytes) > 10000:
                                                    img_pil = Image.open(io.BytesIO(image_bytes))
                                                    img_byte_arr = io.BytesIO()
                                                    img_pil.convert('RGB').save(img_byte_arr, format='JPEG', quality=90)
                                                    images.append({
                                                        'page': None,
                                                        'index': len(images),
                                                        'data': img_byte_arr.getvalue(),
                                                        'ext': 'jpg'
                                                    })
                                        except Exception as e:
                                            logger.warning(f"Failed to extract image from footer: {e}")
                except:
                    pass
        
        full_text = "\n\n".join(text_content)
        logger.info(f"DOCX: Extracted {len(doc.paragraphs)} paragraphs, {len(images)} images, {len(full_text)} chars")
        return full_text, images
        
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return "", []

def extract_xlsx(file_path):
    """Extract text from Excel spreadsheet"""
    text_content = []
    
    try:
        wb = load_workbook(file_path, data_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[SHEET: {sheet_name}]"]
            
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    sheet_text.append(" | ".join(row_text))
            
            if len(sheet_text) > 1:  # Has content beyond header
                text_content.append("\n".join(sheet_text))
        
        full_text = "\n\n".join(text_content)
        logger.info(f"XLSX: Extracted {len(wb.sheetnames)} sheets, {len(full_text)} chars")
        return full_text, []
        
    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        return "", []

def edit_docx(file_path, output_path=None, modifications=None):
    """
    Edit a Word document with various modifications.
    
    Args:
        file_path: Path to the input DOCX file
        output_path: Path to save the modified document (defaults to file_path if None)
        modifications: Dict containing modification instructions:
            - 'add_paragraph': {'text': str, 'style': str (optional)}
            - 'add_heading': {'text': str, 'level': int (1-9)}
            - 'replace_text': {'old': str, 'new': str}
            - 'add_table': {'rows': int, 'cols': int, 'data': list of lists (optional)}
            - 'modify_styles': {'paragraph_font': str, 'paragraph_size': int}
    
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        doc = DocxDocument(file_path)
        
        if modifications is None:
            modifications = {}
        
        # Add new paragraphs
        if 'add_paragraph' in modifications:
            para_data = modifications['add_paragraph']
            if isinstance(para_data, list):
                for item in para_data:
                    p = doc.add_paragraph(item.get('text', ''))
                    if 'style' in item:
                        p.style = item['style']
            else:
                p = doc.add_paragraph(para_data.get('text', ''))
                if 'style' in para_data:
                    p.style = para_data['style']
        
        # Add headings
        if 'add_heading' in modifications:
            heading_data = modifications['add_heading']
            if isinstance(heading_data, list):
                for item in heading_data:
                    doc.add_heading(item.get('text', ''), level=item.get('level', 1))
            else:
                doc.add_heading(heading_data.get('text', ''), level=heading_data.get('level', 1))
        
        # Replace text
        if 'replace_text' in modifications:
            replace_data = modifications['replace_text']
            if isinstance(replace_data, list):
                replacements = replace_data
            else:
                replacements = [replace_data]
            
            for replacement in replacements:
                old_text = replacement.get('old', '')
                new_text = replacement.get('new', '')
                
                # Replace in paragraphs
                for para in doc.paragraphs:
                    if old_text in para.text:
                        for run in para.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
                
                # Replace in tables
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                if old_text in para.text:
                                    for run in para.runs:
                                        if old_text in run.text:
                                            run.text = run.text.replace(old_text, new_text)
        
        # Add tables
        if 'add_table' in modifications:
            table_data = modifications['add_table']
            if isinstance(table_data, list):
                tables = table_data
            else:
                tables = [table_data]
            
            for table_info in tables:
                rows = table_info.get('rows', 2)
                cols = table_info.get('cols', 2)
                table = doc.add_table(rows=rows, cols=cols)
                table.style = 'Table Grid'
                
                # Populate table if data provided
                if 'data' in table_info:
                    data = table_info['data']
                    for i, row_data in enumerate(data[:rows]):
                        for j, cell_data in enumerate(row_data[:cols]):
                            table.rows[i].cells[j].text = str(cell_data)
        
        # Modify styles
        if 'modify_styles' in modifications:
            style_data = modifications['modify_styles']
            
            if 'paragraph_font' in style_data or 'paragraph_size' in style_data:
                for para in doc.paragraphs:
                    for run in para.runs:
                        if 'paragraph_font' in style_data:
                            run.font.name = style_data['paragraph_font']
                        if 'paragraph_size' in style_data:
                            run.font.size = Pt(style_data['paragraph_size'])
        
        # Save the document
        if output_path is None:
            output_path = file_path
        
        doc.save(output_path)
        logger.info(f"DOCX edited successfully: {output_path}")
        return True
        
    except Exception as e:
        logger.error(f"DOCX editing failed: {e}")
        return False

def create_docx(output_path, content=None):
    """
    Create a new Word document from scratch.
    
    Args:
        output_path: Path to save the new document
        content: Dict containing content to add:
            - 'title': str - Document title (Heading 1)
            - 'paragraphs': list of str or dict - Text paragraphs to add
            - 'tables': list of dict - Tables to add
            - 'headings': list of dict - Headings to add
    
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        doc = DocxDocument()
        
        if content is None:
            content = {}
        
        # Add title
        if 'title' in content:
            doc.add_heading(content['title'], level=0)
        
        # Add headings
        if 'headings' in content:
            for heading in content['headings']:
                if isinstance(heading, dict):
                    doc.add_heading(heading.get('text', ''), level=heading.get('level', 1))
                else:
                    doc.add_heading(str(heading), level=1)
        
        # Add paragraphs
        if 'paragraphs' in content:
            for para in content['paragraphs']:
                if isinstance(para, dict):
                    p = doc.add_paragraph(para.get('text', ''))
                    if 'style' in para:
                        p.style = para['style']
                else:
                    doc.add_paragraph(str(para))
        
        # Add tables
        if 'tables' in content:
            for table_info in content['tables']:
                rows = table_info.get('rows', 2)
                cols = table_info.get('cols', 2)
                table = doc.add_table(rows=rows, cols=cols)
                table.style = 'Table Grid'
                
                # Populate table if data provided
                if 'data' in table_info:
                    data = table_info['data']
                    for i, row_data in enumerate(data[:rows]):
                        for j, cell_data in enumerate(row_data[:cols]):
                            table.rows[i].cells[j].text = str(cell_data)
        
        doc.save(output_path)
        logger.info(f"DOCX created successfully: {output_path}")
        return True
        
    except Exception as e:
        logger.error(f"DOCX creation failed: {e}")
        return False
