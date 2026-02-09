import json
import boto3
import logging
import os
import time
import uuid
from urllib.parse import unquote_plus
from decimal import Decimal
from functools import wraps
import contextvars

# -----------------------------------------------------------
# Logging with Correlation ID
# -----------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - [%(correlation_id)s] - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

# Context variable for correlation ID
correlation_id_var = contextvars.ContextVar('correlation_id', default='no-correlation-id')

class CorrelationIdFilter(logging.Filter):
    def filter(self, record):
        record.correlation_id = correlation_id_var.get()
        return True

logger.addFilter(CorrelationIdFilter())

# -----------------------------------------------------------
# Security Utils
# -----------------------------------------------------------
def sanitize_for_logging(text, max_length=100):
    """Sanitize text for logging to prevent log injection attacks."""
    if not text:
        return "[EMPTY]"
    # Convert to string and remove control characters
    text_str = str(text)
    # Remove all control characters including newlines, tabs, etc.
    safe = ''.join(c if c.isprintable() and ord(c) >= 32 else '_' for c in text_str)
    # Truncate and add indicator if truncated
    if len(safe) > max_length:
        return safe[:max_length-3] + "..."
    return safe

def validate_filename(filename):
    if not filename or not isinstance(filename, str):
        raise ValueError("Invalid filename")
    filename = filename.strip()
    if not filename or filename.startswith('.') or '/' in filename or '\\' in filename:
        raise ValueError("Invalid filename")
    if len(filename) > 255:
        raise ValueError("Filename too long")
    return filename

def retry_with_backoff(max_retries=3):
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    if attempt == max_retries - 1:
                        raise
                    time.sleep(2 ** attempt)
        return wrapper
    return decorator

# -----------------------------------------------------------
# AWS Setup
# -----------------------------------------------------------
from botocore.client import Config

# Thread-safe client initialization
_s3_client = None
_bedrock_client = None
_embeddings_client = None
_faiss_cache = {}

def get_s3_client():
    global _s3_client
    if _s3_client is None:
        _s3_client = boto3.client(
            "s3", 
            config=Config(
                signature_version='s3v4', 
                connect_timeout=3, 
                read_timeout=30,
                max_pool_connections=50,
                retries={'max_attempts': 2}
            )
        )
    return _s3_client

# Track when client was created to detect stale connections
_bedrock_client_created_at = None
CLIENT_MAX_AGE_SECONDS = 300  # Refresh client every 5 minutes

def get_bedrock_client():
    """Get Bedrock client with optimized configuration and auto-refresh."""
    global _bedrock_client, _bedrock_client_created_at
    
    # Refresh client if it's too old (prevents stale agent version issues)
    current_time = time.time()
    if _bedrock_client is not None and _bedrock_client_created_at is not None:
        age = current_time - _bedrock_client_created_at
        if age > CLIENT_MAX_AGE_SECONDS:
            logger.info(f"Bedrock client is {age:.0f}s old, refreshing...")
            _bedrock_client = None
    
    if _bedrock_client is None:
        region = os.getenv("AWS_REGION")
        if not region:
            raise ValueError("AWS_REGION must be configured")
        
        _bedrock_client = boto3.client(
            "bedrock-agent-runtime",
            region_name=region,
            config=Config(
                connect_timeout=10, 
                read_timeout=120,
                max_pool_connections=20,
                retries={'max_attempts': 2}
            )
        )
        _bedrock_client_created_at = current_time
        logger.info(f"Bedrock client initialized: {region}")
    return _bedrock_client

def get_embeddings_client():
    """Get embeddings client with proper error handling."""
    global _embeddings_client
    if _embeddings_client is None:
        from langchain_aws import BedrockEmbeddings
        region = os.getenv("AWS_REGION")
        if not region:
            raise ValueError("AWS_REGION must be configured")
        model_id = os.getenv("EMBEDDINGS_MODEL_ID", "cohere.embed-multilingual-v3")
        _embeddings_client = BedrockEmbeddings(model_id=model_id, region_name=region)
        logger.info(f"Embeddings client initialized: {model_id}")
    return _embeddings_client

# Get S3 bucket from environment
BUCKET = os.getenv("S3_BUCKET")
if not BUCKET:
    logger.error("S3_BUCKET environment variable not configured")
    raise ValueError("S3_BUCKET environment variable must be set")

# Track index timestamp to detect updates
_index_s3_timestamp = None

@retry_with_backoff(max_retries=2)
def preload_master_index(force_reload=False):
    """Preload master FAISS index with automatic update detection."""
    from langchain_community.vectorstores import FAISS
    import shutil
    from datetime import datetime
    
    global _index_s3_timestamp
    cache_key = "master_index"
    cache_dir = "/tmp/master_index"
    s3_client = get_s3_client()
    
    # FORCE RELOAD: Always check S3 for updates
    try:
        s3_obj = s3_client.head_object(Bucket=BUCKET, Key="vector_store/master/index.faiss")
        s3_last_modified = s3_obj['LastModified'].timestamp()
        
        # If we have a cached index, check if S3 version is newer
        if cache_key in _faiss_cache and _index_s3_timestamp is not None:
            if s3_last_modified > _index_s3_timestamp:
                logger.info(f"S3 index updated (cached: {datetime.fromtimestamp(_index_s3_timestamp)}, S3: {datetime.fromtimestamp(s3_last_modified)}), reloading...")
                del _faiss_cache[cache_key]
                # Also clear disk cache
                if os.path.exists(cache_dir):
                    shutil.rmtree(cache_dir)
            else:
                logger.info("Master index cached and up-to-date")
                return
        elif cache_key in _faiss_cache:
            logger.info("Master index already cached, checking S3...")
            # Still check if we should reload
            if s3_last_modified > _index_s3_timestamp:
                logger.info("S3 has newer version, reloading...")
                del _faiss_cache[cache_key]
                if os.path.exists(cache_dir):
                    shutil.rmtree(cache_dir)
            else:
                return
            
    except s3_client.exceptions.ClientError as e:
        if e.response.get('Error', {}).get('Code') == '404':
            logger.info("No master index exists yet")
            return
        logger.warning(f"Error checking S3 index timestamp: {type(e).__name__}")
    
    except s3_client.exceptions.ClientError as e:
        if e.response.get('Error', {}).get('Code') == '404':
            logger.info("No master index exists yet")
            return
        logger.warning(f"Error checking S3 index timestamp: {type(e).__name__}")
        return
    
    if os.path.exists(cache_dir):
        shutil.rmtree(cache_dir)
    
    os.makedirs(cache_dir, exist_ok=True)
    
    index_file = os.path.join(cache_dir, "index.faiss")
    pkl_file = os.path.join(cache_dir, "index.pkl")
    
    start_time = time.time()
    # Download files in parallel for faster startup
    import concurrent.futures
    with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
        future1 = executor.submit(s3_client.download_file, BUCKET, "vector_store/master/index.faiss", index_file)
        future2 = executor.submit(s3_client.download_file, BUCKET, "vector_store/master/index.pkl", pkl_file)
        concurrent.futures.wait([future1, future2])
    
    embeddings = get_embeddings_client()
    master_index = FAISS.load_local(cache_dir, embeddings, allow_dangerous_deserialization=True)
    _faiss_cache[cache_key] = master_index
    
    # Store S3 timestamp for future comparisons
    try:
        s3_obj = s3_client.head_object(Bucket=BUCKET, Key="vector_store/master/index.faiss")
        _index_s3_timestamp = s3_obj['LastModified'].timestamp()
    except:
        _index_s3_timestamp = time.time()
    
    logger.info(f"Preloaded master index: {master_index.index.ntotal} vectors in {time.time() - start_time:.2f}s")

try:
    preload_master_index()
except Exception as e:
    logger.error(f"Startup preload failed: {type(e).__name__}")


def list_all_s3_objects(bucket, prefix, max_keys=1000):
    """List all S3 objects with pagination support and limits."""
    if not bucket or not isinstance(bucket, str):
        raise ValueError("Invalid bucket name")
    if not prefix or not isinstance(prefix, str):
        raise ValueError("Invalid prefix")
    
    s3_client = get_s3_client()
    objects = []
    continuation_token = None
    total_fetched = 0
    
    while total_fetched < max_keys:
        params = {
            'Bucket': bucket,
            'Prefix': prefix,
            'MaxKeys': min(1000, max_keys - total_fetched)
        }
        if continuation_token:
            params['ContinuationToken'] = continuation_token
        
        response = s3_client.list_objects_v2(**params)
        batch = response.get('Contents', [])
        objects.extend(batch)
        total_fetched += len(batch)
        
        if not response.get('IsTruncated', False) or total_fetched >= max_keys:
            break
        continuation_token = response.get('NextContinuationToken')
    
    return objects

def batch_delete_s3_objects(bucket, keys):
    """Delete multiple S3 objects in batches of 1000."""
    if not keys:
        return 0
    if not bucket or not isinstance(bucket, str):
        raise ValueError("Invalid bucket name")
    
    s3_client = get_s3_client()
    deleted_count = 0
    
    for i in range(0, len(keys), 1000):
        batch = keys[i:i+1000]
        delete_dict = {'Objects': [{'Key': k} for k in batch], 'Quiet': True}
        try:
            s3_client.delete_objects(Bucket=bucket, Delete=delete_dict)
            deleted_count += len(batch)
        except Exception as e:
            logger.error(f"Batch delete failed: {type(e).__name__}")
            raise
    
    return deleted_count

# -----------------------------------------------------------
# Utilities
# -----------------------------------------------------------
def cors_response(body=None, status=200):
    """Return a standard CORS-enabled response."""
    return {
        "statusCode": status,
        "headers": {
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "GET, POST, DELETE, OPTIONS",
            "Access-Control-Allow-Headers": "Content-Type",
            "Content-Type": "application/json; charset=utf-8",
        },
        "body": json.dumps(body or {}, ensure_ascii=False),
    }


# -----------------------------------------------------------
# /get-upload-url
# -----------------------------------------------------------
def handle_get_upload_url():
    """Generate a presigned S3 URL for direct PDF upload."""
    try:
        timestamp = int(time.time())
        key = f"uploads/upload_{timestamp}.pdf"
        s3_client = get_s3_client()
        url = s3_client.generate_presigned_url(
            "put_object",
            Params={"Bucket": BUCKET, "Key": key, "ContentType": "application/pdf"},
            ExpiresIn=7200,
        )
        logger.info(f"Upload URL generated | Key: {sanitize_for_logging(key)} | Expires: 7200s")
        return cors_response({"uploadUrl": url, "key": key})
    except Exception as e:
        logger.error(f"Upload URL generation failed | Error: {type(e).__name__} | Details: {str(e)[:100]}")
        return cors_response({"error": "Failed to generate upload URL"}, 500)


# -----------------------------------------------------------
# /list-files
# -----------------------------------------------------------
def handle_list_files():
    """List processed PDFs available for querying."""
    try:
        logger.info("Listing processed files | Bucket: {BUCKET} | Prefix: processed/")
        s3_client = get_s3_client()
        response = s3_client.list_objects_v2(Bucket=BUCKET, Prefix="processed/", MaxKeys=1000)
        files = []
        if "Contents" in response:
            for obj in response["Contents"]:
                key = obj["Key"]
                if key.endswith(".json"):
                    filename = key.split("/")[-1].replace(".json", "")
                    files.append({
                        "filename": filename,
                        "lastModified": obj["LastModified"].isoformat(),
                        "size": obj["Size"]
                    })
        
        logger.info(f"List files complete | Count: {len(files)} | Truncated: {response.get('IsTruncated', False)}")
        return cors_response({"files": files})
    except Exception as e:
        logger.error(f"List files failed | Error: {type(e).__name__} | Details: {str(e)[:100]}")
        return cors_response({"error": "Failed to list files"}, 500)


# -----------------------------------------------------------
# /delete-file
# -----------------------------------------------------------
def handle_delete_file(filename):
    """Delete a processed PDF and its associated data."""
    try:
        filename = validate_filename(filename)
        logger.info(f"Delete operation started | File: {sanitize_for_logging(filename)}")
        
        s3_client = get_s3_client()
        
        # Delete processed JSON
        processed_key = f"processed/{filename}.json"
        try:
            s3_client.delete_object(Bucket=BUCKET, Key=processed_key)
            logger.info(f"Deleted processed JSON | Key: {sanitize_for_logging(processed_key)}")
        except Exception as e:
            logger.warning(f"Could not delete processed | Key: {processed_key} | Error: {type(e).__name__}")
        
        # Delete uploads (limit search)
        upload_objects = list_all_s3_objects(BUCKET, f"uploads/{filename}", max_keys=100)
        logger.info(f"Found {len(upload_objects)} upload objects")
        for obj in upload_objects:
            try:
                s3_client.delete_object(Bucket=BUCKET, Key=obj["Key"])
            except Exception as e:
                logger.warning(f"Delete upload failed | Key: {obj['Key']} | Error: {type(e).__name__}")
        
        # Delete images
        image_objects = list_all_s3_objects(BUCKET, f"images/{filename}/", max_keys=500)
        image_keys = [obj["Key"] for obj in image_objects]
        if image_keys:
            deleted_count = batch_delete_s3_objects(BUCKET, image_keys)
            logger.info(f"Deleted images | Count: {deleted_count} | File: {filename}")
        
        # Delete vector store data
        vector_objects = list_all_s3_objects(BUCKET, f"vector_store/{filename}/", max_keys=100)
        vector_keys = [obj["Key"] for obj in vector_objects]
        if vector_keys:
            deleted_count = batch_delete_s3_objects(BUCKET, vector_keys)
            logger.info(f"Deleted vector objects | Count: {deleted_count} | File: {filename}")
        
        # Clear FAISS cache
        if filename in _faiss_cache:
            del _faiss_cache[filename]
            logger.info(f"Cleared FAISS cache | File: {sanitize_for_logging(filename)}")
        
        logger.info(f"Delete operation complete | File: {filename}")
        return cors_response({"message": f"Successfully deleted {filename}"})
    except ValueError as e:
        logger.warning(f"Invalid filename | Error: {type(e).__name__} | Details: {str(e)[:100]}")
        return cors_response({"error": str(e)}, 400)
    except Exception as e:
        logger.error(f"Delete failed | File: {filename} | Error: {type(e).__name__} | Details: {str(e)[:100]}")
        return cors_response({"error": "Failed to delete file"}, 500)


# -----------------------------------------------------------
# Optimized Search Function
# -----------------------------------------------------------
def optimized_search(query, top_k=30):
    """Intelligent search using semantic similarity."""
    try:
        # Use cached index first
        cache_key = "master_index"
        if cache_key in _faiss_cache:
            master_index = _faiss_cache[cache_key]
        else:
            # Fallback: try to load from disk cache
            from langchain_community.vectorstores import FAISS
            cache_dir = "/tmp/master_index"
            
            if os.path.exists(os.path.join(cache_dir, "index.faiss")):
                # Load from existing disk cache
                embeddings = get_embeddings_client()
                master_index = FAISS.load_local(cache_dir, embeddings, allow_dangerous_deserialization=True)
                _faiss_cache[cache_key] = master_index
                
                # Store S3 timestamp
                try:
                    s3_obj = s3_client.head_object(Bucket=BUCKET, Key="vector_store/master/index.faiss")
                    _index_s3_timestamp = s3_obj['LastModified'].timestamp()
                except:
                    _index_s3_timestamp = time.time()
                    
                logger.info("Loaded index from disk cache")
            else:
                # Last resort: download and cache
                preload_master_index()
                if cache_key not in _faiss_cache:
                    logger.error("Failed to load master index")
                    return []
                master_index = _faiss_cache[cache_key]
        
        # Direct semantic search with higher k for better recall
        start_time = time.time()
        results = master_index.similarity_search_with_score(query, k=top_k)
        search_time = time.time() - start_time
        
        # Minimal result formatting for speed
        formatted_results = []
        for doc, semantic_score in results:
            # Extract file information from metadata
            source = doc.metadata.get('source', 'Unknown')
            source_file = doc.metadata.get('source_file', source)
            s3_key = doc.metadata.get('s3_key', '')
            uploaded_name = doc.metadata.get('uploaded_name', source)
            
            # Use the most informative file name available
            file_path = uploaded_name if uploaded_name else (source_file if source_file else source)
            
            formatted_results.append({
                'content': doc.page_content,
                'metadata': doc.metadata,
                'semantic_score': float(semantic_score),
                'source_file': file_path,  # Add explicit file path
                's3_key': s3_key,  # Include S3 key if available
                'document_name': source  # Include base document name
            })
        
        logger.info(f"Semantic search: {len(results)} results in {search_time:.3f}s")
        return formatted_results
        
    except Exception as e:
        logger.error(f"Search failed: {type(e).__name__}")
        return []


# -----------------------------------------------------------
# /search
# -----------------------------------------------------------
def handle_search(query):
    """Handle search requests with intelligent filtering."""
    try:
        if not query or not query.strip():
            return cors_response({"error": "Query is required"}, 400)
        
        query = query.strip()
        logger.info(f"Search started | Query: {sanitize_for_logging(query)} | Length: {len(query)}")
        
        results = optimized_search(query, top_k=15)
        
        logger.info(f"Search complete | Results: {len(results)} | Query: {sanitize_for_logging(query)[:50]}")
        return cors_response({
            "results": results,
            "query": query,
            "total": len(results)
        })
        
    except Exception as e:
        logger.error(f"Search failed | Query: {sanitize_for_logging(query)[:50]} | Error: {type(e).__name__} | Details: {str(e)[:100]}")
        return cors_response({"error": "Search failed"}, 500)


# -----------------------------------------------------------
# Bedrock Agent Integration
# -----------------------------------------------------------
def handle_agent_query(query, session_id=None):
    """Handle queries using Bedrock Agent with search integration."""
    try:
        if not query or not query.strip():
            return cors_response({"error": "Query is required"}, 400)
        
        query = query.strip()
        if not session_id:
            session_id = str(uuid.uuid4())
        
        logger.info(f"Processing agent query: {sanitize_for_logging(query)} (session: {sanitize_for_logging(session_id)})")
        
        # Get agent configuration from environment
        agent_id = os.getenv("BEDROCK_AGENT_ID")
        agent_alias_id = os.getenv("BEDROCK_AGENT_ALIAS_ID")
        
        if not agent_id or not agent_alias_id:
            logger.error("Bedrock agent configuration missing")
            return cors_response({"error": "Agent configuration not found"}, 500)
        
        bedrock_client = get_bedrock_client()
        
        # Invoke the agent
        response = bedrock_client.invoke_agent(
            agentId=agent_id,
            agentAliasId=agent_alias_id,
            sessionId=session_id,
            inputText=query
        )
        
        # Process the response stream
        response_text = ""
        for event in response.get('completion', []):
            if 'chunk' in event:
                chunk = event['chunk']
                if 'bytes' in chunk:
                    response_text += chunk['bytes'].decode('utf-8')
        
        logger.info(f"Agent response generated (length: {len(response_text)})")
        
        return cors_response({
            "response": response_text,
            "sessionId": session_id,
            "query": query
        })
        
    except Exception as e:
        logger.error(f"Agent query failed: {sanitize_for_logging(str(e))}")
        return cors_response({"error": "Agent query failed"}, 500)


# -----------------------------------------------------------
# Lambda Handler
# -----------------------------------------------------------
def lambda_handler(event, context):
    """Main Lambda handler with comprehensive error handling."""
    try:
        logger.info(f"Received event: {sanitize_for_logging(json.dumps(event, default=str))}")
        
        # Handle CORS preflight
        if event.get("httpMethod") == "OPTIONS":
            return cors_response()
        
        # Extract path and method
        path = event.get("path", "")
        method = event.get("httpMethod", "")
        
        # Parse body for POST requests
        body = {}
        if method == "POST" and event.get("body"):
            try:
                body = json.loads(event["body"])
            except json.JSONDecodeError as e:
                logger.warning(f"Invalid JSON in request body: {sanitize_for_logging(str(e))}")
                return cors_response({"error": "Invalid JSON"}, 400)
        
        # Route requests
        if path == "/get-upload-url" and method == "GET":
            return handle_get_upload_url()
        
        elif path == "/list-files" and method == "GET":
            return handle_list_files()
        
        elif path.startswith("/delete-file/") and method == "DELETE":
            filename = unquote_plus(path.split("/delete-file/")[1])
            return handle_delete_file(filename)
        
        elif path == "/search" and method == "POST":
            query = body.get("query", "")
            return handle_search(query)
        
        elif path == "/agent" and method == "POST":
            query = body.get("query", "")
            session_id = body.get("sessionId")
            return handle_agent_query(query, session_id)
        
        else:
            logger.warning(f"Unknown endpoint: {sanitize_for_logging(path)} {sanitize_for_logging(method)}")
            return cors_response({"error": "Endpoint not found"}, 404)
    
    except Exception as e:
        logger.error(f"Unhandled error in lambda_handler: {sanitize_for_logging(str(e))}")
        return cors_response({"error": "Internal server error"}, 500)
        for item in response.get("Contents", []):
            key = item["Key"]
            if key.lower().endswith(".pdf"):
                files.append(os.path.basename(key))
        logger.info(f" Listed {len(files)} processed PDFs.")
        return cors_response({"files": files})
    except Exception as e:
        logger.error(f" Failed to list files: {type(e).__name__}")
        return cors_response({"error": "Failed to list files"}, 500)


# -----------------------------------------------------------
# /agent-query
# -----------------------------------------------------------
def handle_get_upload_url_api(event):
    """API Gateway handler for generating upload URLs."""
    try:
        body_str = event.get("body") or "{}"
        try:
            body = json.loads(body_str) if body_str else {}
        except json.JSONDecodeError:
            return cors_response({"error": "Invalid JSON in request body"}, 400)
        
        file_name = body.get("fileName", f"upload_{int(time.time())}.pdf")
        
        try:
            file_name = validate_filename(file_name)
        except ValueError as e:
            return cors_response({"error": str(e)}, 400)
        
        allowed_extensions = ('.pdf', '.pptx', '.docx', '.xlsx', '.jpg', '.jpeg', '.png', '.tiff')
        if not file_name.lower().endswith(allowed_extensions):
            return cors_response({"error": "Only PDF, PPTX, DOCX, XLSX, JPG, JPEG, PNG, and TIFF files are allowed"}, 400)
        
        # Check if file already exists in uploads/
        s3_client = get_s3_client()
        base_name_no_ext = file_name.rsplit('.', 1)[0] if '.' in file_name else file_name
        try:
            response = s3_client.list_objects_v2(Bucket=BUCKET, Prefix="uploads/")
            if 'Contents' in response:
                for obj in response['Contents']:
                    existing_file = obj['Key'].replace('uploads/', '')
                    existing_no_ext = existing_file.rsplit('.', 1)[0] if '.' in existing_file else existing_file
                    if existing_no_ext == base_name_no_ext:
                        return cors_response({"error": f"File '{file_name}' already exists. Please delete it first or rename your file."}, 409)
        except Exception as e:
            logger.error(f"Error checking duplicates: {type(e).__name__}")
        
        # Ensure filename is properly UTF-8 encoded for S3
        key = f"uploads/{file_name}"
        
        # Determine content type based on file extension
        content_type_map = {
            '.pdf': 'application/pdf',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.tiff': 'image/tiff'
        }
        
        # Use manual split instead of os.path.splitext to preserve Hebrew/Unicode
        file_ext = ('.' + file_name.split('.')[-1]).lower() if '.' in file_name else ''
        content_type = content_type_map.get(file_ext, 'application/octet-stream')
        
        signed_url = s3_client.generate_presigned_url(
            "put_object",
            Params={"Bucket": BUCKET, "Key": key, "ContentType": content_type},
            ExpiresIn=3600
        )
        logger.info(f" Generated upload URL for key: {key}")
        return cors_response({"signedUrl": signed_url, "fileName": file_name})
    except Exception as e:
        logger.error(f" Upload URL error: {type(e).__name__}", exc_info=True)
        return cors_response({"error": "Upload URL generation failed"}, 500)

def handle_list_files_api(event):
    """API Gateway handler for listing files."""
    try:
        import html
        objects = list_all_s3_objects(BUCKET, "processed/")
        filenames = []
        for obj in objects:
            if obj["Key"].endswith(".json"):
                # Use string split instead of os.path.basename to preserve Unicode
                filename = obj["Key"].split("/")[-1]
                # Decode HTML entities
                filename = html.unescape(filename)
                filenames.append(filename)
        return cors_response({"filenames": filenames})
    except Exception as e:
        logger.error(f" List files error: {type(e).__name__}")
        return cors_response({"error": "Failed to list files"}, 500)

def get_document_context(metadata, base_name):
    """Extract rich document-level context for chunk enrichment."""
    context_parts = [f"Document: {base_name}"]
    
    # Add text preview summary
    text_preview = metadata.get('text_preview', '').strip()
    if text_preview:
        first_line = text_preview.split('\n')[0][:80].strip()
        if first_line:
            context_parts.append(f"Subject: {first_line}")
    
    # Add visual context from first image
    images = metadata.get('images', [])
    if images:
        first_img = images[0].get('description', '')
        if 'BRAND:' in first_img:
            brand = first_img.split('BRAND:')[1].split('\n')[0].strip()
            if brand and brand.lower() != 'none visible':
                context_parts.append(f"Brand: {brand}")
    
    return " | ".join(context_parts)

def rebuild_master_index():
    """Rebuild master index from all remaining processed documents."""
    try:
        from langchain_community.vectorstores import FAISS
        from langchain.docstore.document import Document
        import tempfile
        
        logger.info(" Rebuilding master index...")
        
        # Get all processed documents
        processed_objects = list_all_s3_objects(BUCKET, "processed/")
        if not processed_objects:
            logger.info(" No documents to index, deleting master index")
            try:
                s3_client = get_s3_client()
                s3_client.delete_object(Bucket=BUCKET, Key="vector_store/master/index.faiss")
                s3_client.delete_object(Bucket=BUCKET, Key="vector_store/master/index.pkl")
            except:
                pass
            # Clear cache
            if "master_index" in _faiss_cache:
                del _faiss_cache["master_index"]
            return
        
        embeddings = get_embeddings_client()
        all_docs = []
        
        # Collect all documents from processed files
        for obj in processed_objects:
            if not obj['Key'].endswith('.json'):
                continue
            
            try:
                # Get processed metadata
                s3_client = get_s3_client()
                response = s3_client.get_object(Bucket=BUCKET, Key=obj['Key'])
                metadata = json.loads(response['Body'].read().decode('utf-8'))
                
                source_file = metadata.get('source_file', '')
                # Use string split instead of os.path.basename to preserve Hebrew/Unicode
                base_name = source_file.split('/')[-1].replace('.pdf', '')
                
                # Get rich document context
                context_prefix = get_document_context(metadata, base_name)
                
                # Get full text from metadata (not just preview)
                full_text = metadata.get('full_text', metadata.get('text_preview', ''))
                if full_text:
                    # Use smaller chunks for better ranking
                    from langchain.text_splitter import RecursiveCharacterTextSplitter
                    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
                    chunks = splitter.split_text(full_text)
                    
                    for i, chunk in enumerate(chunks):
                        doc = Document(
                            page_content=f"Document: {base_name}\n{chunk}",
                            metadata={
                                "source": base_name,
                                "source_file": source_file,  # Full S3 path
                                "s3_key": source_file,  # Include S3 key
                                "chunk_id": i
                            }
                        )
                        all_docs.append(doc)
                
                # Add image documents with surrounding text context
                images = metadata.get('images', [])
                # Filter out small images BEFORE numbering
                filtered_images = []
                for img_meta in images:
                    s3_key = img_meta.get('s3_key', '')
                    # Skip tiny images (same logic as search)
                    if s3_key:
                        try:
                            s3_obj = s3_client.head_object(Bucket=BUCKET, Key=s3_key)
                            file_size = s3_obj.get('ContentLength', 0)
                            if file_size < 10000:  # 10KB minimum
                                continue
                        except:
                            pass
                    filtered_images.append(img_meta)
                
                # Now number the filtered images
                for idx, img_meta in enumerate(filtered_images, 1):
                    img_desc = img_meta.get('description', 'Image')
                    s3_key = img_meta.get('s3_key', '')
                    page = img_meta.get('page')
                    
                    # Extract text context around the image (few lines before/after)
                    text_context = img_meta.get('text_context', '')
                    if not text_context and full_text and page:
                        # Fallback: extract text near the page number
                        lines = full_text.split('\n')
                        context_lines = []
                        for i, line in enumerate(lines):
                            if f'page {page}' in line.lower() or f'עמוד {page}' in line:
                                # Get 3 lines before and after
                                start = max(0, i-3)
                                end = min(len(lines), i+4)
                                context_lines = lines[start:end]
                                break
                        text_context = ' '.join(context_lines).strip()[:300]
                    
                    # Build page content with context
                    if page is not None:
                        page_content = f"Document: {base_name}\nIMAGE NUMBER {idx} | Image #{idx} | Image number {idx} | תמונה מספר {idx}\nPage {page}: {img_desc}\nContext: {text_context}\nIMAGE_URL:{s3_key}|PAGE:{page}|SOURCE:{base_name}"
                    else:
                        page_content = f"Document: {base_name}\nIMAGE NUMBER {idx} | Image #{idx} | Image number {idx} | תמונה מספר {idx}\n{img_desc}\nContext: {text_context}\nIMAGE_URL:{s3_key}|SOURCE:{base_name}"
                    
                    img_doc = Document(
                        page_content=page_content,
                        metadata={
                            "source": base_name,
                            "source_file": source_file,  # Full S3 path
                            "s3_key": s3_key,
                            "uploaded_name": source_file.split('/')[-1] if source_file else '',
                            "type": "image",
                            "page": page,
                            "image_number": idx,
                            "image_url": img_meta.get('url', ''),
                            "context_summary": context_prefix,
                            "text_context": text_context
                        }
                    )
                    all_docs.append(img_doc)
                    
            except Exception as e:
                logger.error(f" Failed to process {obj['Key']}: {e}")
        
        if not all_docs:
            logger.warning(" No documents collected for indexing")
            return
        
        logger.info(f" Creating master index with {len(all_docs)} documents")
        
        # Create new master index
        with tempfile.TemporaryDirectory() as tmpdir:
            master_store = FAISS.from_documents(all_docs, embeddings)
            master_store.save_local(tmpdir)
            
            # Upload to S3
            s3_client = get_s3_client()
            s3_client.upload_file(os.path.join(tmpdir, "index.faiss"), BUCKET, "vector_store/master/index.faiss")
            s3_client.upload_file(os.path.join(tmpdir, "index.pkl"), BUCKET, "vector_store/master/index.pkl")
        
        # Clear cache to force reload
        if "master_index" in _faiss_cache:
            del _faiss_cache["master_index"]
        
        logger.info(" Master index rebuilt successfully")
        
    except Exception as e:
        logger.error(f" Failed to rebuild master index: {e}")
        raise

def handle_batch_delete_api(event):
    """API Gateway handler for rebuilding index after batch delete."""
    try:
        logger.info("Rebuilding master index after batch deletion")
        rebuild_master_index()
        logger.info("Master index rebuilt successfully")
        return cors_response({"success": True, "message": "Index rebuilt"})
    except Exception as e:
        logger.error(f"Index rebuild error: {type(e).__name__} - {str(e)}")
        return cors_response({"error": f"Index rebuild failed: {str(e)}"}, 500)

def handle_delete_file_api(event):
    """API Gateway handler for deleting files."""
    try:
        from urllib.parse import unquote
        params = event.get("queryStringParameters") or {}
        display_name = unquote(params.get("fileName", ""))
        
        try:
            display_name = validate_filename(display_name)
        except (ValueError, TypeError) as e:
            logger.error(f"Invalid filename: {e}")
            return cors_response({"error": f"Invalid fileName: {str(e)}"}, 400)
        
        logger.info(f"Deleting file: {sanitize_for_logging(display_name)}")
        # Use manual split to preserve Hebrew/Unicode
        base_name = display_name.replace('.json', '') if display_name.endswith('.json') else display_name
        
        # Remove timestamp prefix if present (format: 1234567890_filename)
        if '_' in base_name:
            parts = base_name.split('_', 1)
            if parts[0].isdigit():
                base_name = parts[1]
                logger.info(f"Stripped timestamp prefix, using base_name: {sanitize_for_logging(base_name)}")
        
        # Collect all keys to delete (need to handle both with and without timestamp prefix)
        keys_to_delete = [
            f"processed/{display_name}",  # Original name with timestamp
            f"processed/{base_name}.json",  # Without timestamp
            f"progress/{display_name.replace('.json', '.json')}",  # With timestamp
            f"progress/{base_name}.json"  # Without timestamp
        ]
        
        # Collect uploaded files (original documents)
        try:
            upload_objects = list_all_s3_objects(BUCKET, f"uploads/{base_name}")
            keys_to_delete.extend([obj['Key'] for obj in upload_objects])
            logger.info(f"Found {len(upload_objects)} upload files to delete")
        except Exception as e:
            logger.error(f"Failed to list uploads: {e}")
        
        # Collect images
        try:
            image_objects = list_all_s3_objects(BUCKET, f"images/{base_name}/")
            keys_to_delete.extend([obj['Key'] for obj in image_objects])
            logger.info(f"Found {len(image_objects)} image files to delete")
        except Exception as e:
            logger.error(f"Failed to list images: {e}")
        
        # Batch delete all collected keys
        try:
            deleted_count = batch_delete_s3_objects(BUCKET, keys_to_delete)
            logger.info(f"Deleted {deleted_count} objects for {base_name}")
        except Exception as e:
            logger.error(f"Batch delete failed: {e}")
            return cors_response({"error": f"Delete failed: {str(e)}"}, 500)
        
        # Wait a moment for S3 to propagate deletions
        time.sleep(0.5)
        
        # Rebuild master vector store (removes deleted document from index)
        try:
            logger.info("Rebuilding master index...")
            rebuild_master_index()
            logger.info("Master index rebuilt after deletion")
        except Exception as e:
            logger.error(f"Failed to rebuild master index: {e}")
            # Don't fail the delete operation if index rebuild fails
            logger.warning("Delete succeeded but index rebuild failed - index will be stale")
        
        return cors_response({"success": True, "message": f"{base_name} deleted successfully"})
    except Exception as e:
        logger.error(f"Delete file error: {type(e).__name__} - {str(e)}")
        return cors_response({"error": f"Failed to delete file: {str(e)}"}, 500)

def handle_cancel_upload_api(event):
    """API Gateway handler for cancelling uploads and cleaning up partial files."""
    try:
        params = event.get("queryStringParameters") or {}
        file_name = params.get("fileName")
        
        try:
            file_name = validate_filename(file_name)
        except (ValueError, TypeError) as e:
            return cors_response({"error": "Invalid fileName"}, 400)
        
        logger.info(f" Cancelling upload: {sanitize_for_logging(file_name)}")
        # Use manual split instead of os.path.splitext to preserve Hebrew/Unicode
        base_name = '.'.join(file_name.split('.')[:-1]) if '.' in file_name else file_name
        
        # Create cancellation marker FIRST to stop Lambda processing
        try:
            s3_client = get_s3_client()
            s3_client.put_object(
                Bucket=BUCKET,
                Key=f"cancelled/{base_name}.txt",
                Body=f"Cancelled at {time.time()}",
                ContentType="text/plain"
            )
            logger.info(f" Created cancellation marker")
        except Exception as e:
            logger.error(f" Failed to create cancellation marker: {e}")
        
        # Collect all keys to delete
        keys_to_delete = []
        
        # Collect uploads
        try:
            upload_objects = list_all_s3_objects(BUCKET, f"uploads/{base_name}")
            keys_to_delete.extend([obj['Key'] for obj in upload_objects])
        except Exception as e:
            logger.error(f" Failed to list uploads: {e}")
        
        # Collect processed
        try:
            processed_objects = list_all_s3_objects(BUCKET, f"processed/{base_name}")
            keys_to_delete.extend([obj['Key'] for obj in processed_objects])
        except Exception as e:
            logger.error(f" Failed to list processed: {e}")
        
        # Collect images
        try:
            image_objects = list_all_s3_objects(BUCKET, f"images/{base_name}/")
            keys_to_delete.extend([obj['Key'] for obj in image_objects])
        except Exception as e:
            logger.error(f" Failed to list images: {e}")
        
        # Note: No need to delete vector store - using master index now
        
        # Add cancellation marker to delete list
        keys_to_delete.append(f"cancelled/{base_name}.txt")
        
        # Batch delete all collected keys
        deleted_count = batch_delete_s3_objects(BUCKET, keys_to_delete)
        logger.info(f" Total files deleted: {deleted_count}")
        
        return cors_response({"message": f"Cancelled and deleted {deleted_count} files for {base_name}"})
    except Exception as e:
        logger.error(f" Cancel upload error: {type(e).__name__}")
        return cors_response({"error": "Failed to cancel upload"}, 500)





def handle_search_action(event):
    """Intelligent search using semantic similarity without hardcoded rules."""
    search_start = time.time()
    try:
        logger.info(f"🔍 SEARCH ACTION CALLED - Agent is searching PDFs")
        parse_start = time.time()
        request_body = event.get("requestBody", {})
        content = request_body.get("content", {})
        app_json = content.get("application/json", {})
        properties = app_json.get("properties", [])
        
        query = ""
        for prop in properties:
            if prop.get("name") == "query":
                query = prop.get("value", "")
                break
        
        # Get original user input to check for visual triggers
        original_input = event.get("inputText", "")
        logger.info(f"🔍 ORIGINAL INPUT: {sanitize_for_logging(original_input)}")
        logger.info(f"🔍 SEARCH QUERY: {sanitize_for_logging(query)}")
        logger.info(f"⏱️ SEARCH: Parse request took {time.time() - parse_start:.3f}s")
        
        if not query.strip():
            return {
                "messageVersion": "1.0",
                "response": {
                    "actionGroup": event.get("actionGroup", "LambdaTools"),
                    "apiPath": "/search",
                    "httpMethod": "POST",
                    "httpStatusCode": 200,
                    "responseBody": {"application/json": {"body": json.dumps({"result": "Please specify what you're looking for."})}}
                }
            }
        
        from langchain_community.vectorstores import FAISS
        
        # Use cached index
        cache_start = time.time()
        cache_key = "master_index"
        if cache_key in _faiss_cache:
            master_index = _faiss_cache[cache_key]
            logger.info(f"⏱️ SEARCH: Cache hit took {time.time() - cache_start:.3f}s")
        else:
            logger.info("⏱️ SEARCH: Cache miss, loading index...")
            preload_master_index()
            logger.info(f"⏱️ SEARCH: Index load took {time.time() - cache_start:.3f}s")
            if cache_key not in _faiss_cache:
                return {
                    "messageVersion": "1.0",
                    "response": {
                        "actionGroup": event.get("actionGroup", "LambdaTools"),
                        "apiPath": "/search",
                        "httpMethod": "POST",
                        "httpStatusCode": 200,
                        "responseBody": {"application/json": {"body": json.dumps({"result": "No documents have been uploaded yet. Please upload PDF files first to enable search functionality. You can upload files using the upload feature in the application."})}}
                    }
                }
            master_index = _faiss_cache[cache_key]
        
        # 1. Smart Intent Detection - works for ANY subject
        query_lower = query.lower()
        original_lower = original_input.lower() if original_input else query_lower
        
        # Question words = user wants information (TEXT)
        question_words = ['what', 'how', 'why', 'when', 'where', 'which', 'who', 'whose', 'whom', 'מה', 'איך', 'למה', 'מתי', 'איפה', 'מי']
        # Comparison/analysis = user wants details (TEXT ONLY)
        analysis_words = ['compare', 'difference', 'versus', 'vs', 'between', 'analyze', 'explain', 'describe', 'list', 'summarize', 'השווה', 'הבדל', 'הסבר', 'תאר']
        # Visual commands = user wants to see (IMAGES) - VERY STRONG (English + Hebrew)
        visual_commands = ['show me', 'display', 'picture of', 'photo of', 'image of', 'look at', 'see the', 'view the', 'let me see']
        hebrew_visual_commands = ['תראה', 'הצג', 'תמונה', 'דיאגרמה', 'תרשים', 'דיאגרמות', 'תמונות']
        
        # Count intent signals
        text_signals = 0
        image_signals = 0
        
        # Check for visual commands FIRST (highest priority)
        # English commands (substring match)
        if any(cmd in original_lower for cmd in visual_commands):
            image_signals += 10
            logger.info(f"🖼️ English visual command detected")
        
        # Hebrew commands (word match)
        original_words = original_input.split() if original_input else []
        if any(hebrew_cmd in original_words for hebrew_cmd in hebrew_visual_commands):
            image_signals += 10
            logger.info(f"🖼️ Hebrew visual command detected")
        
        # Check for question words at start (strong text signal)
        first_word = query_lower.split()[0] if query_lower.split() else ''
        if first_word in question_words:
            text_signals += 3
        
        # Check for analysis/comparison words (VERY STRONG text signal)
        if any(word in original_lower for word in analysis_words):
            text_signals += 5
        
        # Check for standalone visual words (weaker signal)
        if any(word in original_lower.split() for word in ['picture', 'photo', 'image', 'drawing', 'תמונה', 'תמונות']):
            image_signals += 1
        
        # Decision based on signal strength
        if image_signals >= 10:
            # Strong visual command detected - IMAGES ONLY
            wants_images = True
            wants_text = False
            logger.info(f"🖼️ Intent: IMAGES (signals: text={text_signals}, image={image_signals})")
        elif text_signals > image_signals:
            wants_text = True
            wants_images = False
            logger.info(f"📄 Intent: TEXT (signals: text={text_signals}, image={image_signals})")
        elif image_signals > text_signals:
            wants_images = True
            wants_text = False
            logger.info(f"🖼️ Intent: IMAGES (signals: text={text_signals}, image={image_signals})")
        else:
            # Equal or no strong signals = hybrid
            wants_text = True
            wants_images = True
            logger.info(f"🔄 Intent: HYBRID (signals: text={text_signals}, image={image_signals})")
        
        # 2. Wide Retrieval with fuzzy matching
        vector_start = time.time()
        
        # Check for specific image number request (e.g., "image 5", "איור מספר 5")
        import re
        image_number_match = re.search(r'(?:image|picture|photo|diagram|תמונה|תמונות|איור|דיאגרמה)\s*(?:number|num|#|מספר)?\s*(\d+)', original_input)
        requested_image_num = None
        if image_number_match:
            requested_image_num = int(image_number_match.group(1))
            logger.info(f"🔢 Specific image number requested: {requested_image_num}")
        
        # Extract document name - try to match against actual document names in index
        doc_filter = None
        
        # For specific image number requests, search MORE results to ensure we find it
        search_k = 100 if requested_image_num else 50
        
        # Extract document name early to use in search
        doc_name_in_query = None
        # Decode HTML entities FIRST (e.g., &quot; -> ")
        import html
        decoded_input = html.unescape(original_input)
        
        # Extended patterns to match various document reference formats
        # Including Hebrew: "מהמסמך" = "from document", 'הצעת מחיר - חברת ליידוס' etc
        patterns = [
            # Quoted filenames (highest priority - exact text in quotes)
            r'[""]([^""]+)["""]',  # Smart quotes
            r'["\'"]([^"\']+)["\']',  # Regular quotes (must be non-greedy and capture full content)
            # Hebrew: "מהמסמך NAME" - capturing until punctuation
            r'מהמסמך\s+["\'"]?([^"\'?]+?)["\'"]?\s*(\?|$)',
            # Hebrew: "ממסמך NAME"
            r'מ(?:ה)?מסמך\s+["\'"]?([^"\'?]+?)["\'"]?\s*(\?|$)',
            # English: "from/in document"
            r'(?:from|in)\s+(?:the\s+)?document[:\s]+["\'"]?([^"\'?]+?)["\'"]?\s*(\?|$)',
        ]
        
        for i, pattern in enumerate(patterns):
            match = re.search(pattern, decoded_input, re.IGNORECASE | re.UNICODE)
            if match:
                # For quoted patterns, group 1 is the content; for others, may be group 1 or 2
                doc_name_in_query = match.group(1).strip()
                logger.info(f"🔍 Pattern {i} matched: '{doc_name_in_query}'")
                # Remove trailing punctuation and file extensions
                doc_name_in_query = doc_name_in_query.rstrip('.,;:!? \t')
                if doc_name_in_query.endswith('.pdf'):
                    doc_name_in_query = doc_name_in_query[:-4]
                logger.info(f"🔍 Extracted doc name: '{doc_name_in_query}'")
                break
        
        # If document name specified, search by document name instead of query
        # This helps narrow down results to the specific document the user is asking about
        if doc_name_in_query:
            search_query = doc_name_in_query
            logger.info(f"🔍 Searching by document name: '{search_query}'")
        else:
            search_query = query
            logger.info(f"🔍 No document name found, searching by query: '{search_query}'")
        
        raw_results = master_index.similarity_search_with_score(search_query, k=search_k)
        logger.info(f"⏱️ SEARCH: Vector search (k={search_k}) took {time.time() - vector_start:.3f}s, found {len(raw_results)} results")
        
        # Get all unique document names from results (try both source and source_file)
        unique_sources = set()
        for doc, _ in raw_results:
            # Add source (base name)
            source = doc.metadata.get('source', '')
            if source:
                unique_sources.add(source)
            # Also add source_file (full path) if available  
            source_file = doc.metadata.get('source_file', '')
            if source_file:
                unique_sources.add(source_file)
            # Also add S3 key if available
            s3_key = doc.metadata.get('s3_key', '')
            if s3_key:
                # Extract just the filename from S3 key
                s3_filename = s3_key.split('/')[-1].replace('.pdf', '')
                if s3_filename:
                    unique_sources.add(s3_filename)
        
        logger.info(f"📁 ALL available documents: {sorted(list(unique_sources))}")
        
        if doc_name_in_query:
            logger.info(f"🔍 Extracted doc name from query: '{doc_name_in_query}'")
        else:
            logger.info("⚠️ No document name extracted from query")
        
        # Match document name - be more flexible with matching
        if doc_name_in_query:
            # Normalize for comparison (handle HTML entities, case, spaces, Hebrew)
            import html
            normalized_query = html.unescape(doc_name_in_query).strip()
            normalized_query_upper = normalized_query.upper()
            
            # Try exact match first (case-insensitive)
            exact_match = None
            for source in unique_sources:
                if normalized_query_upper == html.unescape(source).strip().upper():
                    exact_match = source
                    break
            
            if exact_match:
                doc_filter = exact_match
                logger.info(f"🎯 ✅ Exact match: '{doc_filter}'")
            else:
                # Try substring matching (document name is contained in source)
                # This handles cases where stored name might have prefixes/suffixes
                for source in unique_sources:
                    normalized_source = html.unescape(source).strip().upper()
                    # Remove timestamp prefix if present for comparison
                    if '_' in normalized_source:
                        parts = normalized_source.split('_', 1)
                        if parts[0].isdigit() and len(parts[0]) == 10:  # Unix timestamp
                            normalized_source = parts[1]
                    
                    # Check if query is a substring of source (handles partial names)
                    # Also check reversed to catch cases where source might be part of query
                    if (normalized_query_upper in normalized_source or normalized_source in normalized_query_upper):
                        doc_filter = source
                        logger.info(f"🎯 ✅ Substring match: query='{doc_name_in_query}' matched source='{doc_filter}'")
                        break
                
                # If still no match, try fuzzy matching for typos and partial names
                if not doc_filter:
                    from difflib import SequenceMatcher
                    best_match = None
                    best_ratio = 0
                    for source in unique_sources:
                        # For partial matching, also try removing prefixes
                        source_normalized = html.unescape(source).strip().upper()
                        # Remove timestamp prefix if present (format: 1234567890_filename)
                        if '_' in source_normalized:
                            parts = source_normalized.split('_', 1)
                            if parts[0].isdigit() and len(parts[0]) == 10:  # Unix timestamp
                                source_normalized = parts[1]
                        
                        ratio = SequenceMatcher(None, normalized_query_upper, source_normalized).ratio()
                        if ratio > 0.75 and ratio > best_ratio:  # 75% similarity threshold for better matches
                            best_match = source
                            best_ratio = ratio
                            logger.info(f"🔍 Fuzzy candidate: '{source}' ratio={ratio:.2f}")
                    
                    if best_match:
                        doc_filter = best_match
                        logger.info(f"🎯 ✅ Fuzzy match (ratio={best_ratio:.2f}): query='{doc_name_in_query}' matched source='{doc_filter}'")
                    else:
                        logger.warning(f"⚠️ Document '{doc_name_in_query}' not found in available sources: {sorted(list(unique_sources))}")
            
            # Filter results only if we found a matching document
            if doc_filter:
                # Filter by multiple metadata fields for robustness
                # BACKWARD COMPATIBLE: Falls back to 'source' if new fields don't exist
                filtered = []
                for doc, score in raw_results:
                    source = doc.metadata.get('source', '')
                    source_file = doc.metadata.get('source_file', source)  # Fallback to source
                    s3_key = doc.metadata.get('s3_key', source)  # Fallback to source
                    uploaded_name = doc.metadata.get('uploaded_name', '')
                    
                    # If uploaded_name is empty, try to extract from source
                    if not uploaded_name and source:
                        # Extract filename from path like "1767520048_filename.json"
                        uploaded_name = source.split('/')[-1].replace('.json', '').split('_', 1)[-1] if '_' in source else source
                    
                    # Check if any field matches doc_filter
                    if (source == doc_filter or 
                        source_file == doc_filter or 
                        s3_key.split('/')[-1].replace('.pdf', '') == doc_filter if s3_key else False or
                        uploaded_name == doc_filter):
                        filtered.append((doc, score))
                
                if filtered:
                    raw_results = filtered
                    logger.info(f"🎯 Filtered to {len(raw_results)} results from '{doc_filter}'")
                else:
                    logger.warning(f"⚠️ No results after filtering for '{doc_filter}', using all results")
                    doc_filter = None  # Reset filter to use all results
        
        # Apply fuzzy matching to boost relevant results
        from difflib import SequenceMatcher
        
        def fuzzy_score(text1, text2):
            """Calculate fuzzy similarity between two strings (0-1)."""
            return SequenceMatcher(None, text1.lower(), text2.lower()).ratio()
        
        # Extract key terms from query (ignore stop words)
        stop_words = {'show', 'me', 'the', 'a', 'an', 'in', 'of', 'with', 'from', 'to', 'for', 'and', 'or', 'image', 'picture', 'photo', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'should', 'could', 'may', 'might', 'must', 'can', 'this', 'that', 'these', 'those', 'it', 'its', 'they', 'them', 'their', 'about', 'tell', 'what', 'how', 'why', 'when', 'where', 'file', 'document', 'pdf'}
        query_terms = [w for w in query_lower.split() if w not in stop_words and len(w) > 2]
        
        # Boost results where document name matches query terms
        doc_name_boost = 0
        for term in query_terms:
            if term in ['project', 'chapter', 'section', 'document']:
                continue  # Skip generic terms
            # Check if term appears in any document metadata
            doc_name_boost = 0.3  # Will apply per-result below
        
        # Re-score results with fuzzy matching boost and document name matching
        fuzzy_results = []
        for doc, semantic_score in raw_results:
            content_lower = doc.page_content.lower()
            source = doc.metadata.get('source', '').lower()
            
            # Calculate fuzzy match score for each query term
            fuzzy_boost = 0
            for term in query_terms:
                # VERY STRONG BOOST: Document name matches query term
                if term in source:
                    fuzzy_boost += 5.0
                    logger.info(f"📄 Document name match: '{term}' in '{source}'")
                
                # Check for exact match in content
                if term in content_lower:
                    fuzzy_boost += 1.0
                else:
                    # Check for fuzzy match for typos and variations
                    words_in_content = content_lower.split()
                    best_match = max([fuzzy_score(term, word) for word in words_in_content] + [0])
                    if best_match > 0.75:  # 75% similarity threshold
                        fuzzy_boost += best_match
                        logger.info(f"🔍 Fuzzy match: '{term}' → best_match={best_match:.2f}")
            
            # Adjust semantic score with fuzzy boost (lower score = better)
            adjusted_score = semantic_score - (fuzzy_boost * 0.3)
            fuzzy_results.append((doc, adjusted_score, semantic_score))
        
        # Sort by adjusted score
        fuzzy_results.sort(key=lambda x: x[1])
        raw_results = [(doc, adj_score) for doc, adj_score, _ in fuzzy_results]
        
        # 3. Separate Candidates and filter out unwanted images
        text_candidates = []
        image_candidates = []
        for doc, score in raw_results:
            if doc.metadata.get('type') == 'image':
                desc = doc.page_content.lower()
                s3_key = doc.metadata.get('s3_key', '')
                
                # SKIP filtering if user requested a specific image number
                if not requested_image_num:
                    # Filter out tiny images (logos, icons < 10KB)
                    if s3_key:
                        try:
                            s3_obj = get_s3_client().head_object(Bucket=BUCKET, Key=s3_key)
                            file_size = s3_obj.get('ContentLength', 0)
                            if file_size < 10000:  # 10KB minimum
                                logger.info(f"Filtering out tiny image ({file_size} bytes): {s3_key}")
                                continue
                        except:
                            pass
                    
                    # Filter out ONLY obvious non-technical images (be conservative)
                    exclude_patterns = [
                        ('qr code', 'qr-code', 'barcode'),
                        ('certification badge', 'certificate badge', 'professional badge'),
                        ('company logo', 'brand logo'),
                        ('profile picture', 'avatar', 'headshot')
                    ]
                    
                    should_exclude = False
                    for pattern_group in exclude_patterns:
                        if any(pattern in desc for pattern in pattern_group):
                            logger.info(f"Filtering out: {doc.metadata.get('source', '')} - matched: {pattern_group}")
                            should_exclude = True
                            break
                    
                    if should_exclude:
                        continue
                    
                image_candidates.append((doc, score))
            else:
                text_candidates.append((doc, score))
        
        # 4. Select Best Results based on Intent
        final_results = []
        
        # If specific image number requested, filter to that image only
        if requested_image_num:
            logger.info(f"🔢 Filtering for image #{requested_image_num}, doc_filter='{doc_filter}'")
            
            # Search ALL results (not just image_candidates) for the specific image number
            numbered_images = []
            for doc, score in raw_results:
                if doc.metadata.get('type') != 'image':
                    continue
                    
                img_num = doc.metadata.get('image_number', 0)
                source = doc.metadata.get('source', '')
                
                if doc_filter:
                    # Normalize both for comparison
                    import html
                    normalized_source = html.unescape(source).strip().upper()
                    normalized_filter = html.unescape(doc_filter).strip().upper()
                    
                    if img_num == requested_image_num and normalized_source == normalized_filter:
                        numbered_images.append((doc, score))
                        logger.info(f"✅ MATCH: Image #{img_num} from {source}")
                    elif img_num == requested_image_num:
                        logger.info(f"❌ SKIP: Image #{img_num} from {source} (doc mismatch: '{normalized_source}' != '{normalized_filter}')")
                else:
                    if img_num == requested_image_num:
                        numbered_images.append((doc, score))
                        logger.info(f"✅ MATCH: Image #{img_num} from {source}")
            
            if numbered_images:
                final_results.extend(numbered_images[:1])
                wants_images = True
                wants_text = False
                logger.info(f"🎯 Returning image #{requested_image_num}")
            else:
                # Log all image numbers found for debugging
                all_img_nums = [doc.metadata.get('image_number') for doc, _ in raw_results if doc.metadata.get('type') == 'image']
                logger.warning(f"⚠️ Image #{requested_image_num} not found. Available images: {sorted(set(all_img_nums))}")
                
                if doc_filter:
                    result = f"Image #{requested_image_num} not found in document '{doc_filter}'. Available images: {sorted(set(all_img_nums))}"
                else:
                    result = f"Image #{requested_image_num} not found. Available images: {sorted(set(all_img_nums))}"
                return {
                    "messageVersion": "1.0",
                    "response": {
                        "actionGroup": event.get("actionGroup", "LambdaTools"),
                        "apiPath": "/search",
                        "httpMethod": "POST",
                        "httpStatusCode": 200,
                        "responseBody": {"application/json": {"body": json.dumps({"result": result})}}
                    }
                }
        elif wants_images and not wants_text:
            # PURE IMAGE QUERY
            
            # Check if user is asking for diagrams specifically
            diagram_keywords = ['diagram', 'architecture', 'תרשים', 'ארכיטקטורה', 'דיאגרמה', 'chart', 'flowchart', 'schematic', 'landing zone', 'infrastructure']
            wants_diagrams = any(keyword in original_lower for keyword in diagram_keywords)
            
            # DEBUG: Log metadata for ALL images
            logger.info(f"🔍 DEBUG: Checking {len(image_candidates)} images for diagram_type metadata:")
            for i, (doc, score) in enumerate(image_candidates[:5]):
                img_num = doc.metadata.get('image_number', '?')
                diagram_type = doc.metadata.get('diagram_type')
                logger.info(f"  Image #{img_num}: diagram_type={diagram_type}")
            
            # If user wants diagrams, filter strictly to ONLY diagrams (if any detected)
            if wants_diagrams:
                diagram_candidates = [(doc, score) for doc, score in image_candidates if doc.metadata.get('diagram_type')]
                
                if diagram_candidates:
                    # Found diagrams - ONLY return diagrams, exclude badges/logos
                    logger.info(f"📊 Diagram filter: {len(image_candidates)} images → {len(diagram_candidates)} diagrams (strict)")
                    image_candidates = diagram_candidates
                else:
                    # No diagrams detected - keep all images (detection may have failed)
                    logger.warning(f"⚠️ No diagrams detected in {len(image_candidates)} images - keeping all (detection may have failed)")
                    # Keep original image_candidates unchanged
            
            # Apply attribute filtering for images (colors, sizes, etc.)
            color_words = ['red', 'blue', 'black', 'white', 'gray', 'grey', 'green', 'yellow', 'orange', 'silver', 'brown', 'purple', 'pink', 'gold', 'beige', 'tan']
            size_words = ['large', 'small', 'big', 'tiny', 'huge', 'medium']
            attribute_words = color_words + size_words
            query_attributes = [attr for attr in attribute_words if attr in query_lower]
            
            # Get meaningful query words (not stop words or attributes)
            query_words = [w for w in query_lower.split() if w not in stop_words and w not in attribute_words and len(w) > 2]
            
            scored_images = []
            for doc, semantic_score in image_candidates:
                content_lower = doc.page_content.lower()
                text_context = doc.metadata.get('text_context', '').lower()
                match_score = 0
                
                # CRITICAL: Penalize logos/banners
                is_logo = doc.metadata.get('is_logo_or_banner', False)
                if is_logo:
                    match_score -= 1000
                    logger.info(f"🚫 Logo/banner - penalizing")
                
                # CRITICAL: Prioritize images with diagram_type metadata (from ingestion detection)
                diagram_type = doc.metadata.get('diagram_type', '')
                if diagram_type:
                    match_score += 5000  # MASSIVE boost for detected diagrams (increased from 500)
                    logger.info(f"📊 DIAGRAM DETECTED: {diagram_type}")
                
                # CRITICAL: Match OCR keywords from image (cloud services, technical terms)
                ocr_keywords = doc.metadata.get('ocr_keywords', [])
                img_num = doc.metadata.get('image_number', '?')
                logger.info(f"🔍 Image #{img_num}: ocr_keywords={ocr_keywords}, query_words={query_words}")
                if ocr_keywords and query_words:
                    # Direct keyword matching
                    keyword_matches = sum(1 for kw in ocr_keywords if any(qw in kw.lower() or kw.lower() in qw for qw in query_words))
                    
                    # Semantic matching for AWS-specific terms
                    aws_terms = ['landing', 'zone', 'control', 'tower', 'organization', 'account', 'governance']
                    if any(term in query_lower for term in aws_terms) and 'aws' in [k.lower() for k in ocr_keywords]:
                        keyword_matches += 2
                        logger.info(f"🎯 AWS semantic match: query has AWS terms, image has 'aws' keyword")
                    
                    # Landing zone = VPC/networking architecture
                    if 'landing' in query_lower and 'zone' in query_lower:
                        vpc_keywords = ['vpc', 'subnet', 'network', 'fargate', 'ecs', 'iam']
                        vpc_matches = sum(1 for kw in ocr_keywords if kw.lower() in vpc_keywords)
                        if vpc_matches > 0:
                            keyword_matches += vpc_matches * 2
                            logger.info(f"🎯 Landing zone match: {vpc_matches} VPC/networking keywords found")
                    
                    if keyword_matches > 0:
                        match_score += keyword_matches * 200  # HUGE boost for OCR keyword matches
                        logger.info(f"🔑 OCR keyword matches: {keyword_matches} keywords matched query")
                    else:
                        logger.info(f"❌ No OCR keyword matches for image #{img_num}")
                else:
                    logger.info(f"⚠️ Image #{img_num}: ocr_keywords empty or no query_words")
                
                # STRONG BOOST: If user wants diagrams, prioritize architecture/system diagrams
                if wants_diagrams:
                    if 'architecture diagram' in content_lower or 'system diagram' in content_lower:
                        match_score += 100
                        logger.info(f"📊 Architecture/System diagram detected")
                    elif 'network diagram' in content_lower or 'technical diagram' in content_lower:
                        match_score += 50
                        logger.info(f"📊 Network/Technical diagram detected")
                    elif 'flowchart' in content_lower:
                        match_score += 10
                        logger.info(f"📊 Flowchart detected (lower priority)")
                
                # Boost images that match requested attributes
                if query_attributes:
                    matching_attrs = sum(1 for attr in query_attributes if attr in content_lower)
                    if matching_attrs > 0:
                        match_score += matching_attrs * 30
                    else:
                        match_score -= 20  # Penalize if no attributes match
                
                # CONTEXT MATCHING: Boost images where surrounding text matches query
                if query_words and text_context:
                    context_matches = sum(1 for word in query_words if word in text_context)
                    if context_matches > 0:
                        match_score += context_matches * 20  # Strong boost for context match
                        logger.info(f"🎯 Context match: {context_matches} words in surrounding text")
                
                # Boost images that match subject words in description
                if query_words:
                    word_matches = sum(1 for word in query_words if word in content_lower)
                    match_score += word_matches * 10
                
                # CRITICAL: Use aggressive multiplier to ensure OCR keyword matches dominate ranking
                # Lower combined_score = better ranking (FAISS distance metric)
                combined_score = semantic_score - (match_score * 1.0)  # Increased from 0.1 to 1.0
                scored_images.append((doc, combined_score))
            
            scored_images.sort(key=lambda x: x[1])
            
            # Dynamic image count based on user request
            num_images = 1  # Default: return best match
            
            # Check if user explicitly asks for multiple images
            if 'all' in original_lower and any(word in original_lower for word in ['image', 'diagram', 'picture', 'photo']):
                num_images = min(len(scored_images), 20)  # Return all, max 20
                logger.info(f"User requested ALL images, returning {num_images}")
            else:
                # Check for explicit number request (e.g., "5 diagrams", "show 3 images")
                import re
                num_match = re.search(r'\b(\d+)\s*(?:image|diagram|picture|photo)', original_lower)
                if num_match:
                    num_images = min(int(num_match.group(1)), 20)  # Max 20
                    logger.info(f"User requested {num_images} images")
                else:
                    # Confidence-based: if top result has low confidence, return more options
                    if scored_images:
                        best_score = scored_images[0][1]
                        # If best score is high (weak match), return more options
                        if best_score > 0.8:  # Weak confidence
                            num_images = 9
                            logger.info(f"Low confidence (score={best_score:.2f}), returning {num_images} options")
                        elif best_score > 0.5:  # Medium confidence
                            num_images = 5
                            logger.info(f"Medium confidence (score={best_score:.2f}), returning {num_images} options")
                        else:  # High confidence
                            num_images = 1
                            logger.info(f"High confidence (score={best_score:.2f}), returning {num_images} image")
            
            final_results.extend(scored_images[:num_images])
            
        elif wants_text and not wants_images:
            # PURE TEXT QUERY - return top 5 most relevant chunks with FULL context
            final_results.extend(text_candidates[:5])
            logger.info(f"📄 TEXT ONLY: Selected top {len(text_candidates[:5])} text chunks (full context), skipped {len(image_candidates)} images")
            
        else:
            # HYBRID QUERY - balanced mix with full context
            final_results.extend(text_candidates[:3])
            final_results.extend(image_candidates[:2])
            final_results.sort(key=lambda x: x[1])
        
        format_start = time.time()
        all_results = []
        for doc, score in final_results:
            source = doc.metadata.get('source', 'unknown')
            page = doc.metadata.get('page', '')
            img_num = doc.metadata.get('image_number', '')
            
            # DO NOT strip image URLs. Pass the raw content to the agent.
            raw_content = doc.page_content
            
            if not raw_content:
                continue
            
            # Add image number and page to source for images
            if doc.metadata.get('type') == 'image' and img_num:
                source_label = f"{source} (Image #{img_num}, Page {page})"
            else:
                source_label = source
            
            all_results.append((raw_content, source_label))
        logger.info(f"⏱️ SEARCH: Format results took {time.time() - format_start:.3f}s")
        logger.info(f"⏱️ SEARCH: TOTAL search took {time.time() - vector_start:.3f}s with {len(all_results)} results, best_score={final_results[0][1] if final_results else 'N/A'}")
        logger.info(f"Intent: wants_images={wants_images}, wants_text={wants_text} | Candidates: text={len(text_candidates)}, images={len(image_candidates)} | Final: {len(final_results)}")
        
        # Log first 3 results for debugging
        if final_results:
            for i, (doc, score) in enumerate(final_results[:3]):
                content_preview = doc.page_content[:100].replace('\n', ' ')
                doc_type = doc.metadata.get('type', 'text')
                logger.info(f"Result {i+1}: score={score:.3f}, type={doc_type}, content={content_preview}...")
        
        if not all_results:
            result = "No relevant information found."
        else:
            # [OLD BUGGY CODE REMOVED]
            # if has_images:
            #    ... only returned IMAGE_URLs ...
            # else:
            #    ... returned text ...
            
            # Return results based on intent - NO TRUNCATION, full context
            result_parts = []
            for content, source in all_results:
                clean_content = content.strip()
                # For TEXT-ONLY queries, strip out IMAGE_URL lines
                if wants_text and not wants_images:
                    if 'IMAGE_URL:' not in clean_content:
                        result_parts.append(f"[{source}] {clean_content}")
                else:
                    result_parts.append(f"[{source}] {clean_content}")
                
            result = "\n\n".join(result_parts) if result_parts else "No relevant information found."
        
        # Calculate approximate token usage
        approx_tokens = len(result) // 4  # Rough estimate: 4 chars per token
        approx_cost = (approx_tokens / 1_000_000) * 0.80  # $0.80 per 1M tokens for Haiku
        
        logger.info(f"⏱️ SEARCH: Search returned {len(all_results)} results, result length: {len(result)}")
        logger.info(f"💰 TOKEN USAGE: ~{approx_tokens:,} tokens (~${approx_cost:.4f} cost)")
        logger.info(f"⏱️ SEARCH: TOTAL handle_search_action took {time.time() - search_start:.3f}s")
        logger.info(f"✅ SEARCH COMPLETE - Returning {len(all_results)} results to agent")
        
        # CRITICAL: Bedrock Agent has a response size limit (~10KB for input)
        # If response is too large, truncate it intelligently
        MAX_RESPONSE_SIZE = 6000  # 6KB to be safe (reduced from 8KB)
        if len(result) > MAX_RESPONSE_SIZE:
            logger.warning(f"⚠️ Response too large ({len(result)} bytes), truncating to {MAX_RESPONSE_SIZE} bytes")
            # Truncate but keep complete results
            truncated_parts = []
            current_size = 0
            for content, source in all_results:
                clean_content = content.strip()
                if wants_text and not wants_images and 'IMAGE_URL:' in clean_content:
                    continue
                result_text = f"[{source}] {clean_content}"
                if current_size + len(result_text) > MAX_RESPONSE_SIZE:
                    break
                truncated_parts.append(result_text)
                current_size += len(result_text) + 2
            result = "\n\n".join(truncated_parts) + "\n\n[Some results omitted due to size limit]"
        
        response_body = json.dumps({"result": result})
        logger.info(f"📦 Response body size: {len(response_body)} bytes")
        
        return {
            "messageVersion": "1.0",
            "response": {
                "actionGroup": event.get("actionGroup", "LambdaTools"),
                "apiPath": "/search",
                "httpMethod": "POST",
                "httpStatusCode": 200,
                "responseBody": {
                    "application/json": {
                        "body": response_body
                    }
                }
            }
        }
    except Exception as e:
        import traceback
        error_msg = str(e)
        stack_trace = traceback.format_exc()
        logger.error(f"❌ SEARCH FAILED: {type(e).__name__}")
        logger.error(f"Error details: {error_msg}")
        logger.error(f"Stack trace: {stack_trace}")
        
        # Return error in proper format
        return {
            "messageVersion": "1.0",
            "response": {
                "actionGroup": event.get("actionGroup", "LambdaTools"),
                "apiPath": "/search",
                "httpMethod": "POST",
                "httpStatusCode": 200,
                "responseBody": {
                    "application/json": {
                        "body": json.dumps({"result": f"I encountered an error while searching. Please try again. Error: {type(e).__name__}"})
                    }
                }
            }
        }

# SES Email Action - DISABLED to save $64/month on NAT Gateway
# SES requires internet access (no VPC endpoint), which needs NAT Gateway
# Uncomment this function and restore NAT Gateway if email functionality is needed

# def handle_send_email_action(event):
#     """Send email via SES - called by Bedrock Agent."""
#     try:
#         request_body = event.get("requestBody", {})
#         content = request_body.get("content", {})
#         app_json = content.get("application/json", {})
#         properties = app_json.get("properties", [])
#         
#         to_email = ""
#         subject = ""
#         body = ""
#         
#         for prop in properties:
#             name = prop.get("name", "")
#             value = prop.get("value", "")
#             if name == "to_email":
#                 to_email = value
#             elif name == "subject":
#                 subject = value
#             elif name == "body":
#                 body = value
#         
#         logger.info(f"Sending email to: {sanitize_for_logging(to_email)}")
#         
#         if not to_email or not subject or not body:
#             return {
#                 "messageVersion": "1.0",
#                 "response": {
#                     "actionGroup": event.get("actionGroup", "LambdaTools"),
#                     "apiPath": "/send-email",
#                     "httpMethod": "POST",
#                     "httpStatusCode": 400,
#                     "responseBody": {
#                         "application/json": {
#                             "body": json.dumps({"error": "Missing required fields: to_email, subject, body"})
#                         }
#                     }
#                 }
#             }
#         
#         # Send email via SES
#         region = os.getenv("AWS_REGION")
#         if not region:
#             raise ValueError("AWS_REGION must be configured")
#         
#         sender_email = os.getenv("SES_SENDER_EMAIL")
#         if not sender_email:
#             raise ValueError("SES_SENDER_EMAIL must be configured")
#         
#         ses = boto3.client('ses', region_name=region)
#         
#         response = ses.send_email(
#             Source=sender_email,
#             Destination={'ToAddresses': [to_email]},
#             Message={
#                 'Subject': {'Data': subject, 'Charset': 'UTF-8'},
#                 'Body': {'Text': {'Data': body, 'Charset': 'UTF-8'}}
#             }
#         )
#         
#         logger.info(f" Email sent successfully. MessageId: {response['MessageId']}")
#         
#         return {
#             "messageVersion": "1.0",
#             "response": {
#                 "actionGroup": event.get("actionGroup", "LambdaTools"),
#                 "apiPath": "/send-email",
#                 "httpMethod": "POST",
#                 "httpStatusCode": 200,
#                 "responseBody": {
#                     "application/json": {
#                         "body": json.dumps({"result": f"Email sent successfully to {to_email}"})
#                     }
#                 }
#             }
#         }
#     
#     except Exception as e:
#         logger.error(f"Email send failed: {sanitize_for_logging(str(e))}")
#         return {
#             "messageVersion": "1.0",
#             "response": {
#                 "actionGroup": event.get("actionGroup", "LambdaTools"),
#                 "apiPath": "/send-email",
#                 "httpMethod": "POST",
#                 "httpStatusCode": 500,
#                 "responseBody": {
#                     "application/json": {
#                         "body": json.dumps({"error": "Email sending failed"})
#                     }
#                 }
#             }
#         }





def handle_agent_query(event):
    """Route query to Bedrock Agent for natural language response."""
    total_start = time.time()
    try:
        parse_start = time.time()
        body_str = event.get("body", "{}")
        try:
            body = json.loads(body_str) if body_str else {}
        except json.JSONDecodeError as e:
            logger.error(f"Invalid JSON | Error: {type(e).__name__} | Details: {str(e)[:100]}")
            return cors_response({"error": "Invalid JSON in request body"}, 400)
        
        query = body.get("query", "")
        session_id = body.get("sessionId", f"session-{int(time.time())}")
        logger.info(f"⏱️ TIMING: Parse body took {time.time() - parse_start:.3f}s")
        
        # Validate inputs
        if not query or not isinstance(query, str):
            return cors_response({"error": "Query is required and must be a string"}, 400)
        
        if len(query) > 10000:
            return cors_response({"error": "Query too long (max 10000 characters)"}, 400)
        
        config_start = time.time()
        agent_id = os.getenv("BEDROCK_AGENT_ID")
        alias_id = os.getenv("BEDROCK_AGENT_ALIAS_ID")
        
        if not agent_id or not alias_id:
            logger.error("BEDROCK_AGENT_ID or BEDROCK_AGENT_ALIAS_ID environment variable not set")
            return cors_response({"error": "Agent configuration error"}, 500)
        
        try:
            bedrock_agent_runtime = get_bedrock_client()
            logger.info(f"⏱️ TIMING: Get Bedrock client took {time.time() - config_start:.3f}s")
        except Exception as e:
            logger.error(f"Failed to get Bedrock client: {sanitize_for_logging(str(e))}")
            return cors_response({"error": "Failed to initialize agent client"}, 500)
        
        invoke_start = time.time()
        logger.info(f"⏱️ TIMING: Starting agent invocation at {invoke_start - total_start:.3f}s")
        logger.info(f"🤖 AGENT QUERY: {sanitize_for_logging(query)}")
        logger.info(f"📋 Agent ID: {agent_id[:20]}... | Alias: {alias_id[:20]}...")
        
        # Retry logic for rate limiting
        max_retries = 3
        retry_delay = 1
        
        for attempt in range(max_retries):
            try:
                response = bedrock_agent_runtime.invoke_agent(
                    agentId=agent_id,
                    agentAliasId=alias_id,
                    sessionId=session_id,
                    inputText=query,
                    enableTrace=True
                )
                logger.info(f"⏱️ TIMING: Agent invocation completed in {time.time() - invoke_start:.3f}s")
                logger.info(f"✅ Agent response object received, processing stream...")
                break
            except Exception as e:
                error_str = str(e)
                if 'ThrottlingException' in error_str or 'rate' in error_str.lower():
                    if attempt < max_retries - 1:
                        logger.warning(f"Rate limited, retrying in {retry_delay}s (attempt {attempt + 1}/{max_retries})")
                        time.sleep(retry_delay)
                        retry_delay *= 2
                        continue
                logger.error(f"Agent invocation failed after {time.time() - invoke_start:.3f}s: {sanitize_for_logging(str(e))}")
                return cors_response({"error": "Service temporarily busy, please try again"}, 503)
        
        stream_start = time.time()
        answer = ""
        chunk_count = 0
        trace_info = []
        try:
            event_stream = response.get('completion')
            if not event_stream:
                return cors_response({"error": "Invalid agent response"}, 500)
            
            for event in event_stream:
                # Log trace events to see if search is being called
                if 'trace' in event:
                    trace = event['trace']
                    if 'orchestrationTrace' in trace:
                        orch = trace['orchestrationTrace']
                        if 'invocationInput' in orch:
                            logger.info(f"🔍 Agent is invoking action: {json.dumps(orch['invocationInput'])[:200]}")
                        if 'observation' in orch:
                            logger.info(f"📊 Agent received observation: {json.dumps(orch['observation'])[:200]}")
                
                if 'chunk' in event and 'bytes' in event['chunk']:
                    answer += event['chunk']['bytes'].decode('utf-8')
                    chunk_count += 1
            logger.info(f"⏱️ TIMING: Stream processing took {time.time() - stream_start:.3f}s ({chunk_count} chunks, {len(answer)} bytes)")
        except Exception as e:
            full_error = str(e)
            logger.error(f"❌ Failed to process response stream: {full_error}")
            logger.error(f"❌ Answer collected so far (length={len(answer)}): {answer[:200]}")
            return cors_response({"error": "Failed to process agent response"}, 500)
        
        # Extract IMAGE_URL: markers and generate presigned URLs
        image_start = time.time()
        images = []
        import re
        import html
        
        # DEBUG: Log what agent actually returned
        logger.info(f"📝 AGENT RESPONSE (length={len(answer)}): {answer[:500] if len(answer) > 0 else '[EMPTY]'}")
        
        # Check if user explicitly requested images
        query_lower = query.lower()
        visual_keywords = ['show me', 'display', 'picture', 'photo', 'image', 'diagram', 'תראה', 'הצג', 'תמונה', 'דיאגרמה']
        user_wants_images = any(keyword in query_lower for keyword in visual_keywords)
        
        if not user_wants_images:
            logger.info("🚫 User did not request images, skipping image extraction")
            # Strip IMAGE_URL markers from text response
            answer = re.sub(r'IMAGE_URL:[^\n]+\n?', '', answer)
            answer = re.sub(r'images/[^\n]+\.(?:jpg|jpeg|png|gif)[^\n]*\n?', '', answer)
            answer = re.sub(r'\n{3,}', '\n\n', answer).strip()
        else:
            logger.info("✅ User requested images, extracting IMAGE_URLs")
        
        if user_wants_images:
            # Extract image paths - handle multiple formats:
            # Format 1: IMAGE_URL:path|PAGE:n|SOURCE:name (with metadata)
            # Format 2: IMAGE_URL:path (simple format)  
            # Format 3: images/path/file.jpg|SOURCE:name (plain path with metadata)
            # Format 4: images/path/file.jpg (plain S3 path)
            matches = []
            
            # Try IMAGE_URL: prefix format first
            image_url_pattern = r'IMAGE_URL:([^\n|]+)'
            matches = re.findall(image_url_pattern, answer)
            
            # If no IMAGE_URL: prefix, look for plain S3 paths with |SOURCE: metadata
            if not matches:
                plain_with_source = r'(images/[^|\n]+\.(?:jpg|jpeg|png|gif))\|SOURCE:'
                matches = re.findall(plain_with_source, answer)
            
            # If still no matches, look for any S3 image paths
            if not matches:
                plain_path_pattern = r'(images/[^\s\n|]+\.(?:jpg|jpeg|png|gif))'
                matches = re.findall(plain_path_pattern, answer)
            
            logger.info(f"🔍 Found {len(matches)} image paths in response")
            
            if matches:
                logger.info(f"Found {len(matches)} IMAGE_URL markers in agent response")
                s3_client = get_s3_client()
                unique_keys = set()
                failed_images = []
            
            for s3_key in matches:
                # Decode HTML entities (e.g., &quot; -> ")
                s3_key = html.unescape(s3_key.strip())
                # Verify it's a valid image key and not a duplicate
                if s3_key and s3_key not in unique_keys and s3_key.startswith('images/'):
                    unique_keys.add(s3_key)
                    # First verify the image exists in S3
                    try:
                        s3_client.head_object(Bucket=BUCKET, Key=s3_key)
                        logger.info(f"✅ Image exists in S3: {s3_key}")
                    except Exception as head_error:
                        if 'NoSuchKey' in str(type(head_error)) or '404' in str(head_error):
                            logger.error(f"❌ Image NOT found in S3 with key: {s3_key}")
                            # Fallback logic: Search for the image by filename across all 'images/'
                            
                            img_filename = s3_key.split('/')[-1] # Get the filename, e.g., "doc_img15.jpg"
                            
                            if not img_filename:
                                logger.error("❌ Could not extract filename from key. Aborting search.")
                                continue

                            logger.info(f"🔍 FALLBACK: Searching for filename '{img_filename}' in 'images/' prefix...")
                            
                            try:
                                # Use the robust paginated list function
                                all_image_objects = list_all_s3_objects(BUCKET, prefix="images/", max_keys=5000)
                                
                                # Log first 5 filenames for debugging
                                logger.info(f"🔍 First 5 S3 files:")
                                for i, obj in enumerate(all_image_objects[:5]):
                                    logger.info(f"  {i+1}. {obj['Key']}")
                                
                                found = False
                                # Try exact match first
                                for obj in all_image_objects:
                                    obj_key = obj['Key']
                                    if obj_key.endswith(img_filename):
                                        logger.info(f"🎯 FALLBACK SUCCESS (exact): {obj_key}")
                                        s3_key = obj_key
                                        found = True
                                        break
                                
                                # If not found, try matching just the image number (e.g., _img15.jpg)
                                if not found and '_img' in img_filename:
                                    img_num_part = img_filename.split('_img')[-1]  # e.g., "15.jpg"
                                    for obj in all_image_objects:
                                        obj_key = obj['Key']
                                        if f'_img{img_num_part}' in obj_key:
                                            logger.info(f"🎯 FALLBACK SUCCESS (partial): {obj_key}")
                                            s3_key = obj_key
                                            found = True
                                            break
                                
                                if not found:
                                    logger.error(f"❌ FALLBACK FAILED: No match for '{img_filename}' in {len(all_image_objects)} objects")
                                    failed_images.append(img_filename)
                                    continue

                            except Exception as search_e:
                                logger.error(f"❌ FALLBACK FAILED: Error during s3 list: {search_e}")
                                failed_images.append(img_filename)
                                continue # Skip this image
                        else:
                            logger.error(f"Failed to verify image: {head_error}")
                            failed_images.append(s3_key.split('/')[-1] if '/' in s3_key else s3_key)
                            continue
                    
                    # Generate presigned URL
                    try:
                        url = s3_client.generate_presigned_url(
                            'get_object',
                            Params={'Bucket': BUCKET, 'Key': s3_key},
                            ExpiresIn=3600
                        )
                        images.append(url)
                        logger.info(f"Generated presigned URL for: {s3_key}")
                    except Exception as e:
                        logger.error(f"Failed to generate URL for {s3_key}: {e}")
            
            # Extract source info from IMAGE_URL lines before removing them
            source_info = []
            if images:
                for match in re.finditer(r'IMAGE_URL:([^|]+)\|PAGE:(\d+)\|SOURCE:([^\n]+)', answer):
                    s3_key, page, source = match.groups()
                    # Extract image number from s3_key if present
                    img_num_match = re.search(r'_img(\d+)\.', s3_key)
                    img_num = int(img_num_match.group(1)) + 1 if img_num_match else ''
                    if img_num:
                        source_info.append(f"Image #{img_num} from {source} (Page {page})")
                    else:
                        source_info.append(f"Image from {source} (Page {page})")
            
            # If no images were successfully generated, show error
            if not images and failed_images:
                answer = f"I found {len(failed_images)} image(s) but couldn't load them from storage. This may be due to file path encoding issues. Please contact support.\n\nFailed images: {', '.join(failed_images[:3])}"
            else:
                # ALWAYS remove the markers from the text response
                answer = re.sub(r'IMAGE_URL:[^\n]+\n?', '', answer)
                answer = re.sub(r'images/[^\n]+\.(?:jpg|jpeg|png|gif)[^\n]*\n?', '', answer)
                
                # Remove filler phrases
                filler_phrases = [
                    r'Here are the relevant diagrams from the search results:',
                    r'Here are the relevant diagrams for the [^:]+:',
                    r'Here are the images you asked for:',
                    r'The search results contain the following images:',
                ]
                for phrase in filler_phrases:
                    answer = re.sub(phrase, '', answer, flags=re.IGNORECASE)
                
                answer = re.sub(r'\n{3,}', '\n\n', answer).strip()
                
                # CRITICAL: Bedrock Agent requires non-empty response
                # Use a single space if answer is empty but we have images
                if not answer and images:
                    answer = " "
                
                # Add source information only if there's actual text content
                if source_info and answer and answer != " ":
                    answer += "\n\nSource: " + ", ".join(source_info)
            
        logger.info(f"⏱️ TIMING: Image URL generation took {time.time() - image_start:.3f}s ({len(images)} images)")
        
        # Calculate total token usage (input + output)
        input_tokens = len(query) // 4
        output_tokens = len(answer) // 4
        total_tokens = input_tokens + output_tokens
        total_cost = (input_tokens / 1_000_000) * 0.80 + (output_tokens / 1_000_000) * 4.00
        
        logger.info(f"💰 AGENT TOKEN USAGE: Input: ~{input_tokens:,} | Output: ~{output_tokens:,} | Total: ~{total_tokens:,} tokens (~${total_cost:.4f})")
        logger.info(f"⏱️ TIMING: TOTAL agent query took {time.time() - total_start:.3f}s")
        
        return cors_response({"response": answer, "sessionId": session_id, "images": images})
    
    except json.JSONDecodeError as e:
        logger.error(f"JSON decode error | Error: {type(e).__name__} | Details: {str(e)[:100]}")
        return cors_response({"error": "Invalid request format"}, 400)
    except Exception as e:
        error_msg = str(e)
        logger.error(f"Agent query failed | Error: {type(e).__name__} | Details: {error_msg[:200]} | Time: {time.time() - total_start:.2f}s")
        
        # Check if error is related to missing index
        if "index" in error_msg.lower() or "faiss" in error_msg.lower() or "vector" in error_msg.lower():
            return cors_response({"error": "No documents available. Please upload PDF files first to start asking questions."}, 200)
        
        return cors_response({"error": "Failed to process query. Please try again."}, 500)



def handle_get_image(event):
    """Generate presigned URL for image retrieval."""
    try:
        params = event.get("queryStringParameters") or {}
        image_key = params.get("key")
        
        logger.info(f"Get image request for key: {sanitize_for_logging(image_key)}")
        
        if not image_key or not isinstance(image_key, str):
            return cors_response({"error": "Missing or invalid image key"}, 400)
        
        # Validate key starts with images/ to prevent unauthorized access
        if not image_key.startswith("images/"):
            logger.warning(f"Attempted access to non-image key: {sanitize_for_logging(image_key)}")
            return cors_response({"error": "Invalid image key"}, 403)
        
        s3_client = get_s3_client()
        try:
            s3_client.head_object(Bucket=BUCKET, Key=image_key)
            logger.info(f"Image exists: {image_key}")
        except s3_client.exceptions.NoSuchKey:
            logger.error(f"Image not found: {image_key}")
            return cors_response({"error": "Image not found"}, 404)
        except Exception as e:
            logger.error(f"Failed to check image: {type(e).__name__}")
            return cors_response({"error": "Failed to verify image"}, 500)
        
        try:
            url = s3_client.generate_presigned_url(
                "get_object",
                Params={"Bucket": BUCKET, "Key": image_key},
                ExpiresIn=300
            )
            logger.info(f"Generated presigned URL: {image_key}")
            return cors_response({"url": url})
        except Exception as e:
            logger.error(f"Failed to generate URL: {type(e).__name__}")
            return cors_response({"error": "Failed to generate image URL"}, 500)
            
    except Exception as e:
        logger.error(f"Get image error: {sanitize_for_logging(str(e))}")
        return cors_response({"error": "Failed to get image"}, 500)

def handle_processing_status(event):
    """Check processing status of uploaded file."""
    try:
        params = event.get("queryStringParameters") or {}
        file_name = params.get("fileName")
        
        try:
            file_name = validate_filename(file_name)
        except (ValueError, TypeError) as e:
            return cors_response({"error": "Invalid fileName"}, 400)
        
        base_name = os.path.splitext(file_name)[0] if '.' in file_name else file_name
        logger.info(f" Checking status for file: {sanitize_for_logging(file_name)}")
        
        s3_client = get_s3_client()
        
        # Check for errors FIRST
        try:
            error_obj = s3_client.get_object(Bucket=BUCKET, Key=f"errors/{base_name}.txt")
            error_msg = error_obj['Body'].read().decode('utf-8')
            return cors_response({
                "status": "failed",
                "progress": 0,
                "message": f"Processing failed: {error_msg}"
            })
        except s3_client.exceptions.NoSuchKey:
            pass
        except Exception as e:
            logger.error(f"Error checking error marker: {type(e).__name__}")
        
        # Check if cancelled
        try:
            s3_client.head_object(Bucket=BUCKET, Key=f"cancelled/{base_name}.txt")
            return cors_response({
                "status": "cancelled",
                "progress": 0,
                "message": "Processing was cancelled"
            })
        except s3_client.exceptions.NoSuchKey:
            pass
        except s3_client.exceptions.ClientError as e:
            if e.response.get('Error', {}).get('Code') != '404':
                logger.error(f"Error checking cancellation: {type(e).__name__}")
        except Exception as e:
            logger.error(f"Error checking cancellation: {type(e).__name__}")
        
        # Check for progress updates FIRST (more current than processed marker)
        try:
            progress_obj = s3_client.get_object(Bucket=BUCKET, Key=f"progress/{base_name}.json")
            progress_data = json.loads(progress_obj['Body'].read().decode('utf-8'))
            logger.info(f"Found progress marker: {progress_data}")
            return cors_response({
                "status": progress_data.get("status", "processing"),
                "progress": progress_data.get("progress", 50),
                "message": progress_data.get("message", "Processing...")
            })
        except s3_client.exceptions.NoSuchKey:
            logger.info(f"Progress marker not found")
        except Exception as e:
            logger.error(f"Error checking progress: {type(e).__name__}")
        
        # Check if processing is complete
        try:
            logger.info(f"Checking: s3://{BUCKET}/processed/{base_name}.json")
            processed_obj = s3_client.get_object(Bucket=BUCKET, Key=f"processed/{base_name}.json")
            try:
                processed_data = json.loads(processed_obj['Body'].read().decode('utf-8'))
            except (json.JSONDecodeError, KeyError, TypeError) as json_err:
                logger.error(f"Failed to parse processed marker: {json_err}")
                # File exists but is corrupted - treat as completed anyway
                return cors_response({
                    "status": "completed",
                    "progress": 100,
                    "message": "Processing complete"
                })
            logger.info(f"Found completion marker: {base_name}")
            return cors_response({
                "status": "completed",
                "progress": 100,
                "message": "Processing complete",
                "data": processed_data
            })
        except s3_client.exceptions.NoSuchKey:
            logger.info(f"Processed marker not found")
        except Exception as e:
            logger.error(f"Error checking processed: {type(e).__name__}")
        
        # Check if file exists in uploads (fallback)
        try:
            s3_client.head_object(Bucket=BUCKET, Key=f"uploads/{file_name}")
            logger.info(f"File in uploads, processing")
            return cors_response({
                "status": "processing",
                "progress": 50,
                "message": "Processing document..."
            })
        except s3_client.exceptions.NoSuchKey:
            pass
        except Exception as e:
            logger.error(f"Error checking uploads: {type(e).__name__}")
        
        # File not found anywhere
        return cors_response({
            "status": "not_found",
            "progress": 0,
            "message": "File not found or processing not started"
        })
        
    except Exception as e:
        logger.error(f" Processing status error: {type(e).__name__}")
        return cors_response({"error": "Failed to check processing status"}, 500)


# -----------------------------------------------------------
# WebSocket Authorizer
# -----------------------------------------------------------
def handle_websocket_authorizer(event):
    """Validate Cognito JWT token for WebSocket connections."""
    try:
        import jwt
        import requests
        from jwt.algorithms import RSAAlgorithm
        
        # Extract token from query string
        query_params = event.get('queryStringParameters', {})
        token = query_params.get('token') if query_params else None
        
        if not token:
            logger.warning("No token provided in WebSocket connection")
            return generate_auth_policy('user', 'Deny', event['methodArn'])
        
        # Get Cognito configuration
        region = os.getenv('AWS_REGION', 'us-east-1')
        user_pool_id = os.getenv('USER_POOL_ID')
        client_id = os.getenv('CLIENT_ID')
        
        if not user_pool_id or not client_id:
            logger.error("Cognito configuration missing")
            return generate_auth_policy('user', 'Deny', event['methodArn'])
        
        # Fetch JWKS keys
        jwks_url = f'https://cognito-idp.{region}.amazonaws.com/{user_pool_id}/.well-known/jwks.json'
        jwks = requests.get(jwks_url, timeout=5).json()
        
        # Decode header to get kid
        unverified_header = jwt.get_unverified_header(token)
        kid = unverified_header['kid']
        
        # Find the key
        key = None
        for jwk_key in jwks['keys']:
            if jwk_key['kid'] == kid:
                key = RSAAlgorithm.from_jwk(json.dumps(jwk_key))
                break
        
        if not key:
            logger.warning("Key not found in JWKS")
            return generate_auth_policy('user', 'Deny', event['methodArn'])
        
        # Verify token
        decoded = jwt.decode(
            token,
            key,
            algorithms=['RS256'],
            audience=client_id,
            options={'verify_exp': True}
        )
        
        user_id = decoded.get('sub', 'unknown')
        logger.info(f"Token verified for user: {user_id}")
        
        return generate_auth_policy(user_id, 'Allow', event['methodArn'], {
            'userId': user_id,
            'email': decoded.get('email', ''),
            'username': decoded.get('cognito:username', '')
        })
        
    except jwt.ExpiredSignatureError:
        logger.warning("Token expired")
        return generate_auth_policy('user', 'Deny', event['methodArn'])
    except jwt.InvalidTokenError as e:
        logger.warning(f"Invalid token: {e}")
        return generate_auth_policy('user', 'Deny', event['methodArn'])
    except Exception as e:
        logger.error(f"Authorization failed: {e}")
        return generate_auth_policy('user', 'Deny', event['methodArn'])

def generate_auth_policy(principal_id, effect, resource, context=None):
    """Generate IAM policy for API Gateway."""
    auth_response = {
        'principalId': principal_id,
        'policyDocument': {
            'Version': '2012-10-17',
            'Statement': [{
                'Action': 'execute-api:Invoke',
                'Effect': effect,
                'Resource': resource
            }]
        }
    }
    
    if context:
        auth_response['context'] = context
    
    return auth_response

# -----------------------------------------------------------
# WebSocket Support
# -----------------------------------------------------------
def send_websocket_message(connection_url, connection_id, message):
    """Send message to WebSocket connection."""
    try:
        logger.info(f"WS SEND: Creating client for {connection_url}")
        client = boto3.client('apigatewaymanagementapi', endpoint_url=connection_url)
        logger.info(f"WS SEND: Posting to connection {connection_id[:10]}...")
        client.post_to_connection(
            ConnectionId=connection_id,
            Data=json.dumps(message).encode('utf-8')
        )
        logger.info(f"WS SEND: Message sent successfully")
        return True
    except Exception as e:
        logger.error(f"WS SEND FAILED: {type(e).__name__}: {str(e)}")
        import traceback
        logger.error(f"WS SEND TRACE: {traceback.format_exc()}")
        return False

def handle_websocket_agent_query(event):
    """Handle agent query via WebSocket with streaming support."""
    connection_id = event.get('connectionId')
    connection_url = event.get('connectionUrl')
    body_str = event.get('body', '{}')
    
    try:
        logger.info(f"WS: connectionId={connection_id}, url={connection_url}")
        
        # Handle both string and dict body
        if isinstance(body_str, str):
            body = json.loads(body_str)
        else:
            body = body_str
        
        query = body.get('query', '')
        session_id = body.get('sessionId', f"session-{int(time.time())}")
        
        logger.info("WS: Parsed query: {sanitize_for_logging(query)}")
        
        if not query:
            logger.error("WS: No query provided")
            send_websocket_message(connection_url, connection_id, {
                'type': 'error',
                'message': 'Query is required'
            })
            return
        
        logger.info("WS: Sending status update...")
        send_websocket_message(connection_url, connection_id, {
            'type': 'status',
            'message': 'Thinking and searching documents...'
        })
        
        # Get agent configuration
        agent_id = os.getenv("BEDROCK_AGENT_ID")
        agent_alias_id = os.getenv("BEDROCK_AGENT_ALIAS_ID")
        
        logger.info(f"WS: Agent config: {agent_id[:10]}... / {agent_alias_id[:10]}...")
        
        if not agent_id or not agent_alias_id:
            logger.error("WS: Agent config missing")
            send_websocket_message(connection_url, connection_id, {
                'type': 'error',
                'message': 'Agent configuration not found'
            })
            return
        
        logger.info("WS: Getting Bedrock client...")
        bedrock_client = get_bedrock_client()
        
        logger.info("WS: Invoking agent with streaming...")
        response = bedrock_client.invoke_agent(
            agentId=agent_id,
            agentAliasId=agent_alias_id,
            sessionId=session_id,
            inputText=query,
            enableTrace=False  # Disable trace for faster streaming
        )
        
        logger.info("WS: Processing stream in real-time...")
        answer = ""
        chunk_count = 0
        last_send_time = time.time()
        buffer = ""
        
        for event_item in response.get('completion', []):
            if 'chunk' in event_item and 'bytes' in event_item['chunk']:
                chunk_text = event_item['chunk']['bytes'].decode('utf-8')
                answer += chunk_text
                buffer += chunk_text
                chunk_count += 1
                
                # Send immediately when buffer has complete sentences or every 0.1s
                current_time = time.time()
                has_sentence_end = any(char in buffer for char in ['.', '!', '?', '\n'])
                time_elapsed = current_time - last_send_time
                
                if has_sentence_end or time_elapsed > 0.1 or len(buffer) > 100:
                    if buffer.strip():
                        success = send_websocket_message(connection_url, connection_id, {
                            'type': 'chunk',
                            'data': buffer
                        })
                        if not success:
                            logger.warning(f"Failed to send chunk {chunk_count}")
                            break
                        buffer = ""
                        last_send_time = current_time
        
        # Send any remaining buffer
        if buffer.strip():
            send_websocket_message(connection_url, connection_id, {
                'type': 'chunk',
                'data': buffer
            })
        
        logger.info(f"WS: Streamed {chunk_count} chunks, total {len(answer)} bytes")
        
        # Extract images
        images = []
        query_lower = query.lower()
        visual_keywords = ['show me', 'display', 'picture', 'photo', 'image', 'diagram', 'תראה', 'הצג', 'תמונה', 'דיאגרמה']
        user_wants_images = any(keyword in query_lower for keyword in visual_keywords)
        
        if user_wants_images:
            import re
            import html
            image_url_pattern = r'IMAGE_URL:([^\n|]+)'
            matches = re.findall(image_url_pattern, answer)
            
            if matches:
                s3_client = get_s3_client()
                for s3_key in matches:
                    s3_key = html.unescape(s3_key.strip())
                    if s3_key and s3_key.startswith('images/'):
                        try:
                            url = s3_client.generate_presigned_url(
                                'get_object',
                                Params={'Bucket': BUCKET, 'Key': s3_key},
                                ExpiresIn=3600
                            )
                            images.append(url)
                        except Exception as e:
                            logger.error(f"Failed to generate URL: {e}")
            
            answer = re.sub(r'IMAGE_URL:[^\n]+\n?', '', answer)
            answer = re.sub(r'\n{3,}', '\n\n', answer).strip()
        
        logger.info("WS: Sending completion...")
        send_websocket_message(connection_url, connection_id, {
            'type': 'complete',
            'response': answer,
            'images': images,
            'sessionId': session_id
        })
        logger.info("WS: Query complete")
        
    except Exception as e:
        import traceback
        logger.error(f"WS: Query failed: {e}")
        logger.error(f"WS: Traceback: {traceback.format_exc()}")
        try:
            send_websocket_message(connection_url, connection_id, {
                'type': 'error',
                'message': str(e)
            })
        except Exception as send_err:
            logger.error(f"WS: Failed to send error: {send_err}")

# -----------------------------------------------------------
# Lambda Entrypoint
# -----------------------------------------------------------
def lambda_handler(event, context):
    """Main Lambda handler for API Gateway and Bedrock Agent requests."""
    # Set correlation ID for request tracking
    request_id = context.request_id if hasattr(context, 'request_id') else str(uuid.uuid4())
    correlation_id_var.set(request_id)
    
    logger.info(f"Lambda triggered | RequestId: {request_id} | Event: {json.dumps(event)[:500]}")
    
    # Handle WebSocket Authorizer (REQUEST type)
    if event.get('type') == 'REQUEST' and event.get('methodArn'):
        return handle_websocket_authorizer(event)
    
    # Handle WebSocket invocation (detect by requestContext.routeKey)
    request_context = event.get('requestContext', {})
    route_key = request_context.get('routeKey')
    
    if route_key in ['$connect', '$disconnect', 'query']:
        logger.info(f"WebSocket route detected: {route_key}")
        if route_key == '$connect':
            return {'statusCode': 200}
        elif route_key == '$disconnect':
            return {'statusCode': 200}
        elif route_key == 'query':
            try:
                connection_id = request_context.get('connectionId')
                domain_name = request_context.get('domainName')
                stage = request_context.get('stage')
                connection_url = f"https://{domain_name}/{stage}"
                
                logger.info(f"WS HANDLER: Starting query handler")
                logger.info(f"WS HANDLER: connId={connection_id[:10]}...")
                logger.info(f"WS HANDLER: url={connection_url}")
                
                # Pass body directly (already a string from API Gateway)
                ws_event = {
                    'websocket': True,
                    'connectionId': connection_id,
                    'connectionUrl': connection_url,
                    'body': event.get('body', '{}')
                }
                logger.info("WS HANDLER: Calling handle_websocket_agent_query...")
                handle_websocket_agent_query(ws_event)
                logger.info("WS HANDLER: Returned from handle_websocket_agent_query")
                return {'statusCode': 200}
            except Exception as e:
                logger.error(f"WS HANDLER ERROR: {e}")
                import traceback
                logger.error(f"WS HANDLER TRACE: {traceback.format_exc()}")
                return {'statusCode': 500}
    
    # Handle async index rebuild
    if event.get("action") == "rebuild_index":
        logger.info("Starting async index rebuild")
        try:
            rebuild_master_index()
            logger.info("Async index rebuild completed successfully")
            return {"statusCode": 200, "body": json.dumps({"status": "rebuild_complete"})}
        except Exception as e:
            logger.error(f"Async index rebuild failed: {e}")
            return {"statusCode": 500, "body": json.dumps({"status": "rebuild_failed", "error": str(e)})}
    
    # Handle warmup ping from EventBridge
    if event.get("source") == "aws.events" and event.get("detail-type") == "Scheduled Event":
        logger.info("Warmup ping received")
        return {"statusCode": 200, "body": json.dumps({"status": "warm"})}
    
    # Handle async agent query processing
    if event.get("action") == "process_agent_query":
        query_id = event.get("queryId")
        query = event.get("query")
        session_id = event.get("sessionId")
        process_agent_query_background(query_id, query, session_id)
        return {"statusCode": 200}
    
    # Check if this is a Bedrock Agent action invocation
    if "messageVersion" in event and "agent" in event:
        api_path = event.get("apiPath", "")
        logger.info(f"Bedrock Agent action | Path: {api_path}")
        if api_path == "/search":
            return handle_search_action(event)
        # elif api_path == "/send-email":  # DISABLED - SES email action commented out
        #     return handle_send_email_action(event)
        else:
            return {
                "messageVersion": "1.0",
                "response": {
                    "actionGroup": event.get("actionGroup", ""),
                    "apiPath": api_path,
                    "httpMethod": event.get("httpMethod", "POST"),
                    "httpStatusCode": 404,
                    "responseBody": {
                        "application/json": {
                            "body": json.dumps({"error": "Unknown action"})
                        }
                    }
                }
            }
    
    # API Gateway invocations
    path = event.get("path", "")
    method = event.get("httpMethod", "")
    logger.info(f"API Gateway request | Method: {method} | Path: {path}")

    if path.startswith("/production/"):
        path = "/" + path[12:]
    elif path.startswith("/default/"):
        path = "/" + path[8:]
    elif path.startswith("/prod/"):
        path = "/" + path[5:]
    
    if path != event.get("path", ""):
        logger.info(f"Path normalized | Original: {event.get('path')} | Cleaned: {path}")

    if method == "OPTIONS":
        return cors_response()

    try:
        if path == "/upload" and method == "POST":
            return handle_get_upload_url_api(event)
        elif path == "/get-upload-url" and method in ["GET", "POST"]:
            return handle_get_upload_url_api(event)
        elif path == "/list-files" and method == "GET":
            return handle_list_files_api(event)
        elif path == "/delete-file" and method == "DELETE":
            return handle_delete_file_api(event)
        elif path.startswith("/delete-file/") and method == "DELETE":
            filename = unquote_plus(path.split("/delete-file/")[1])
            return handle_delete_file(filename)
        elif path == "/batch-delete" and method == "POST":
            return handle_batch_delete_api(event)
        elif path == "/cancel-upload" and method == "DELETE":
            return handle_cancel_upload_api(event)
        elif path == "/agent-query" and method == "POST":
            return handle_agent_query_async(event)
        elif path == "/agent-status" and method == "GET":
            return handle_agent_status(event)
        elif path == "/get-image" and method == "GET":
            return handle_get_image(event)
        elif path == "/processing-status" and method == "GET":
            return handle_processing_status(event)
        elif path == "/view-file" and method == "GET":
            return handle_view_file(event)
        elif path == "/autofill/extract-source" and method == "POST":
            return handle_autofill_extract_source(event)
        elif path == "/autofill/match-fields" and method == "POST":
            return handle_autofill_match_fields(event)
        elif path == "/autofill/fill-document" and method == "POST":
            return handle_autofill_fill_document(event)
        else:
            logger.error(f"Unhandled route | Path: {path} | Method: {method}")
            return cors_response({"error": f"Unhandled route: {path}"}, 404)

    except Exception as e:
        logger.error(f"Lambda error | Type: {type(e).__name__} | Message: {str(e)[:200]}")
        return cors_response({"error": "Internal server error"}, 500)


# -----------------------------------------------------------
# Async Agent Query (replaces WebSocket)
# -----------------------------------------------------------
def handle_agent_query_async(event):
    """Start async agent query - returns immediately with queryId."""
    try:
        body_str = event.get("body", "{}")
        body = json.loads(body_str) if body_str else {}
        
        query = body.get("query", "")
        session_id = body.get("sessionId", f"session-{int(time.time())}")
        
        if not query or len(query) > 10000:
            return cors_response({"error": "Invalid query"}, 400)
        
        # Generate unique query ID
        query_id = f"query_{int(time.time())}_{uuid.uuid4().hex[:8]}"
        
        # Write initial status to S3
        s3_client = get_s3_client()
        s3_client.put_object(
            Bucket=BUCKET,
            Key=f"agent-status/{query_id}.json",
            Body=json.dumps({"status": "processing", "query": query, "sessionId": session_id}),
            ContentType="application/json"
        )
        
        logger.info(f"Query {query_id} started")
        
        # Invoke Lambda async to process query
        lambda_client = boto3.client('lambda')
        lambda_client.invoke(
            FunctionName=os.getenv('AWS_LAMBDA_FUNCTION_NAME'),
            InvocationType='Event',
            Payload=json.dumps({
                "action": "process_agent_query",
                "queryId": query_id,
                "query": query,
                "sessionId": session_id
            })
        )
        
        return cors_response({"queryId": query_id, "status": "processing"})
        
    except Exception as e:
        logger.error(f"Agent query failed: {e}")
        return cors_response({"error": "Failed to start query"}, 500)


def process_agent_query_background(query_id, query, session_id):
    """Process agent query in background and write result to S3."""
    s3_client = get_s3_client()
    status_key = f"agent-status/{query_id}.json"
    
    try:
        agent_id = os.getenv("BEDROCK_AGENT_ID")
        alias_id = os.getenv("BEDROCK_AGENT_ALIAS_ID")
        
        if not agent_id or not alias_id:
            raise ValueError("Agent config missing")
        
        bedrock_client = get_bedrock_client()
        
        # Update status: Searching
        s3_client.put_object(
            Bucket=BUCKET,
            Key=status_key,
            Body=json.dumps({"status": "processing", "message": "Searching documents..."}),
            ContentType="application/json"
        )
        time.sleep(0.5)  # Give frontend time to poll
        
        # Invoke agent with retry
        for attempt in range(3):
            try:
                response = bedrock_client.invoke_agent(
                    agentId=agent_id,
                    agentAliasId=alias_id,
                    sessionId=session_id,
                    inputText=query,
                    enableTrace=False
                )
                break
            except Exception as e:
                if 'Throttling' in str(e) and attempt < 2:
                    time.sleep(2 ** attempt)
                    continue
                raise
        
        
        # Update status: Generating
        s3_client.put_object(
            Bucket=BUCKET,
            Key=status_key,
            Body=json.dumps({"status": "processing", "message": "Generating response..."}),
            ContentType="application/json"
        )
        
        # Process response
        answer = ""
        for event in response.get('completion', []):
            if 'chunk' in event and 'bytes' in event['chunk']:
                answer += event['chunk']['bytes'].decode('utf-8')
        
        # Extract images - ALWAYS check for images to complement answers
        images = []
        import re, html
        matches = re.findall(r'IMAGE_URL:([^\n|]+)', answer)
        for s3_key in matches:
            s3_key = html.unescape(s3_key.strip())
            if s3_key.startswith('images/'):
                try:
                    url = s3_client.generate_presigned_url(
                        'get_object',
                        Params={'Bucket': BUCKET, 'Key': s3_key},
                        ExpiresIn=3600
                    )
                    images.append(url)
                except:
                    pass
        answer = re.sub(r'IMAGE_URL:[^\n]+\n?', '', answer).strip()
        
        # Write success to S3
        s3_client.put_object(
            Bucket=BUCKET,
            Key=status_key,
            Body=json.dumps({
                "status": "completed",
                "response": answer,
                "images": images,
                "sessionId": session_id
            }),
            ContentType="application/json"
        )
        
        logger.info(f"Query {query_id} completed")
        
    except Exception as e:
        logger.error(f"Query {query_id} failed: {e}")
        try:
            s3_client.put_object(
                Bucket=BUCKET,
                Key=status_key,
                Body=json.dumps({"status": "failed", "error": str(e)}),
                ContentType="application/json"
            )
        except:
            pass


def handle_view_file(event):
    """Generate presigned URL for viewing uploaded file."""
    try:
        params = event.get("queryStringParameters") or {}
        file_name = params.get("fileName")
        
        if not file_name:
            return cors_response({"error": "fileName required"}, 400)
        
        logger.info(f"View file request: {sanitize_for_logging(file_name)}")
        
        s3_client = get_s3_client()
        
        # Remove .json extension (from processed/ folder listing)
        base_name = file_name.replace('.json', '')
        
        # Remove timestamp prefix if present (format: 1234567890_filename)
        if '_' in base_name:
            parts = base_name.split('_', 1)
            if parts[0].isdigit() and len(parts[0]) == 10:
                base_name = parts[1]
                logger.info(f"Stripped timestamp: {base_name}")
        
        # List ALL files in uploads/ and find exact match (ignoring extension)
        try:
            response = s3_client.list_objects_v2(Bucket=BUCKET, Prefix="uploads/")
            if 'Contents' not in response:
                return cors_response({"error": "No files in uploads"}, 404)
            
            # Search for file with matching name (any extension)
            for obj in response['Contents']:
                s3_filename = obj['Key'].replace('uploads/', '')
                # Remove extension from both for comparison
                s3_name_no_ext = s3_filename.rsplit('.', 1)[0] if '.' in s3_filename else s3_filename
                base_name_no_ext = base_name.rsplit('.', 1)[0] if '.' in base_name else base_name
                
                if s3_name_no_ext == base_name_no_ext:
                    s3_key = obj['Key']
                    logger.info(f"Found match: {s3_key}")
                    
                    # Generate URL
                    url = s3_client.generate_presigned_url(
                        'get_object',
                        Params={'Bucket': BUCKET, 'Key': s3_key},
                        ExpiresIn=300
                    )
                    return cors_response({"url": url})
            
            logger.error(f"No match for: {base_name}")
            return cors_response({"error": "File not found"}, 404)
            
        except Exception as e:
            logger.error(f"Error: {type(e).__name__}")
            return cors_response({"error": "File not found"}, 404)
        
    except Exception as e:
        logger.error(f"View error: {type(e).__name__}")
        return cors_response({"error": "Failed to generate view URL"}, 500)


def handle_agent_status(event):
    """Check status of agent query."""
    try:
        params = event.get("queryStringParameters") or {}
        query_id = params.get("queryId")
        
        if not query_id:
            return cors_response({"error": "queryId required"}, 400)
        
        s3_client = get_s3_client()
        status_key = f"agent-status/{query_id}.json"
        
        try:
            response = s3_client.get_object(Bucket=BUCKET, Key=status_key)
            status_data = json.loads(response['Body'].read().decode('utf-8'))
            
            # Delete status file after returning completed result
            if status_data.get('status') in ['completed', 'failed']:
                try:
                    s3_client.delete_object(Bucket=BUCKET, Key=status_key)
                except:
                    pass
            
            return cors_response(status_data)
            
        except s3_client.exceptions.NoSuchKey:
            return cors_response({"status": "not_found"}, 404)
            
    except Exception as e:
        logger.error(f"Status check failed: {e}")
        return cors_response({"error": "Failed to check status"}, 500)




# ===================================================================
# DOCUMENT AUTO-FILL HANDLERS
# ===================================================================

def handle_autofill_extract_source(event):
    """Save source document and trigger parsing via ingestion Lambda."""
    try:
        logger.info("Autofill extract source called")
        body_str = event.get("body", "{}")
        body = json.loads(body_str) if body_str else {}
        
        session_id = body.get("sessionId", f"autofill_{int(time.time())}")
        check_status = body.get("checkStatus", False)
        
        # Status check mode
        if check_status:
            s3_client = get_s3_client()
            text_key = f"document-autofill/sessions/{session_id}/source_text.txt"
            try:
                s3_client.head_object(Bucket=BUCKET, Key=text_key)
                return cors_response({"status": "ready"})
            except:
                return cors_response({"status": "processing"})
        
        # Upload mode
        file_data = body.get("fileData")
        filename = body.get("filename")
        
        logger.info(f"Session: {session_id}, Filename: {filename}, Has data: {bool(file_data)}")
        
        if not file_data or not filename:
            return cors_response({"error": "Missing fileData or filename"}, 400)
        
        import base64
        try:
            file_bytes = base64.b64decode(file_data)
            logger.info(f"Decoded {len(file_bytes)} bytes")
        except Exception as decode_err:
            logger.error(f"Base64 decode failed: {decode_err}")
            return cors_response({"error": "Invalid file data"}, 400)
        
        s3_client = get_s3_client()
        source_key = f"document-autofill/sessions/{session_id}/source_{filename}"
        s3_client.put_object(Bucket=BUCKET, Key=source_key, Body=file_bytes)
        logger.info(f"Saved to S3: {source_key}")
        
        # Invoke ingestion Lambda to parse document
        lambda_client = boto3.client('lambda')
        ingestion_function = os.getenv('INGESTION_LAMBDA_NAME', 'pdfquery-ingestion-worker')
        logger.info(f"Invoking {ingestion_function}")
        
        lambda_client.invoke(
            FunctionName=ingestion_function,
            InvocationType='Event',
            Payload=json.dumps({
                "action": "parse_autofill_document",
                "sessionId": session_id,
                "s3Key": source_key,
                "filename": filename
            })
        )
        
        logger.info(f"Triggered parsing for {filename}")
        return cors_response({"sessionId": session_id, "status": "processing"})
        
    except Exception as e:
        logger.error(f"Extract source failed: {e}", exc_info=True)
        return cors_response({"error": str(e)}, 500)



def handle_autofill_match_fields(event):
    """Save target document and trigger parsing."""
    try:
        body_str = event.get("body", "{}")
        body = json.loads(body_str) if body_str else {}
        
        session_id = body.get("sessionId")
        file_data = body.get("fileData")
        filename = body.get("filename")
        
        if not session_id or not file_data or not filename:
            return cors_response({"error": "Missing required fields"}, 400)
        
        import base64
        file_bytes = base64.b64decode(file_data)
        
        s3_client = get_s3_client()
        target_key = f"document-autofill/sessions/{session_id}/target_{filename}"
        s3_client.put_object(Bucket=BUCKET, Key=target_key, Body=file_bytes)
        
        # Trigger parsing for target document
        lambda_client = boto3.client('lambda')
        ingestion_function = os.getenv('INGESTION_LAMBDA_NAME', 'pdfquery-ingestion-worker')
        target_text_key = f"document-autofill/sessions/{session_id}/target_text.txt"
        lambda_client.invoke(
            FunctionName=ingestion_function,
            InvocationType='Event',
            Payload=json.dumps({
                "action": "parse_autofill_document",
                "sessionId": session_id,
                "s3Key": target_key,
                "filename": filename,
                "outputKey": target_text_key
            })
        )
        
        logger.info(f"Saved and triggered parsing for target: {filename}")
        return cors_response({"status": "ready"})
        
    except Exception as e:
        logger.error(f"Match fields failed: {e}", exc_info=True)
        return cors_response({"error": str(e)}, 500)



def handle_autofill_fill_document(event):
    """Fill target document fields in-place using LLM to map source data."""
    try:
        body_str = event.get("body", "{}")
        body = json.loads(body_str) if body_str else {}
        session_id = body.get("sessionId")
        if not session_id:
            return cors_response({"error": "Missing sessionId"}, 400)
        s3_client = get_s3_client()
        
        # Get source text
        text_key = f"document-autofill/sessions/{session_id}/source_text.txt"
        try:
            text_obj = s3_client.get_object(Bucket=BUCKET, Key=text_key)
            source_text = text_obj['Body'].read().decode('utf-8')
        except s3_client.exceptions.NoSuchKey:
            return cors_response({"error": "Source not ready. Please wait and try again."}, 404)
        
        # Get target document
        target_objects = s3_client.list_objects_v2(Bucket=BUCKET, Prefix=f"document-autofill/sessions/{session_id}/target_")
        if not target_objects.get('Contents'):
            return cors_response({"error": "Target document not found"}, 404)
        target_key = target_objects['Contents'][0]['Key']
        target_filename = target_key.split('/')[-1].replace('target_', '')
        file_ext = target_filename.rsplit('.', 1)[-1].lower() if '.' in target_filename else 'txt'
        
        # Download original target file
        import io
        target_obj = s3_client.get_object(Bucket=BUCKET, Key=target_key)
        target_bytes = io.BytesIO(target_obj['Body'].read())
        
        # Handle TXT files with simple LLM filling
        if file_ext == 'txt':
            target_text = target_bytes.read().decode('utf-8')
            bedrock = boto3.client('bedrock-runtime', region_name=os.getenv("AWS_REGION"))
            prompt = f"""Fill this form using ONLY exact values from source. Do NOT calculate.\n\nSOURCE:\n{source_text}\n\nFORM:\n{target_text}\n\nRULES:\n1. Copy exact values only\n2. Match yearly to yearly, monthly to monthly\n3. If not in source, write "Data not available"\n4. NO calculations\n5. Keep Question:/Answer: format\n\nFilled form:"""
            response = bedrock.invoke_model(modelId="us.anthropic.claude-sonnet-4-5-20250929-v1:0", body=json.dumps({"anthropic_version": "bedrock-2023-05-31", "max_tokens": 2000, "messages": [{"role": "user", "content": prompt}]}))
            result = json.loads(response['body'].read())
            filled_text = result['content'][0]['text']
            output_filename = target_filename.rsplit('.', 1)[0] + '_filled.txt'
            filled_key = f"document-autofill/completed/{session_id}_{output_filename}"
            s3_client.put_object(Bucket=BUCKET, Key=filled_key, Body=filled_text.encode('utf-8'), ContentType='text/plain')
            download_url = s3_client.generate_presigned_url('get_object', Params={'Bucket': BUCKET, 'Key': filled_key}, ExpiresIn=3600)
            logger.info(f"TXT document filled: {filled_key}")
            return cors_response({"downloadUrl": download_url, "filename": output_filename})
        
        # Extract fields from target document
        fields = []
        if file_ext == 'docx':
            from docx import Document
            doc = Document(target_bytes)
            for para in doc.paragraphs:
                if para.text.strip():
                    fields.append({'type': 'paragraph', 'text': para.text, 'obj': para})
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            fields.append({'type': 'cell', 'text': cell.text, 'obj': cell})
        elif file_ext == 'xlsx':
            from openpyxl import load_workbook
            wb = load_workbook(target_bytes)
            ws = wb.active
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value:
                        fields.append({'type': 'cell', 'text': str(cell.value), 'obj': cell})
        elif file_ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(target_bytes)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        fields.append({'type': 'shape', 'text': shape.text, 'obj': shape})
        else:
            return cors_response({"error": "Unsupported file format"}, 400)
        
        if not fields:
            return cors_response({"error": "No fillable fields found in target"}, 400)
        
        # Use LLM to map source data to each field
        bedrock = boto3.client('bedrock-runtime', region_name=os.getenv("AWS_REGION"))
        field_mappings = []
        
        for field in fields:
            prompt = f"""Given this source data and a field from a form, determine the appropriate value to fill.

SOURCE DATA:
{source_text}

FIELD TEXT: {field['text']}

If this field should be filled with data from the source, respond with ONLY the exact value to fill (no explanations).
If this is a label/header or no matching data exists, respond with: SKIP

Value:"""
            
            response = bedrock.invoke_model(
                modelId="us.anthropic.claude-sonnet-4-5-20250929-v1:0",
                body=json.dumps({
                    "anthropic_version": "bedrock-2023-05-31",
                    "max_tokens": 200,
                    "messages": [{"role": "user", "content": prompt}]
                })
            )
            result = json.loads(response['body'].read())
            value = result['content'][0]['text'].strip()
            
            if value and value != "SKIP":
                field_mappings.append({'field': field, 'value': value})
        
        # Fill fields in original document
        if file_ext == 'docx':
            for mapping in field_mappings:
                obj = mapping['field']['obj']
                if mapping['field']['type'] == 'paragraph':
                    obj.text = mapping['value']
                elif mapping['field']['type'] == 'cell':
                    obj.text = mapping['value']
            output_buffer = io.BytesIO()
            doc.save(output_buffer)
            file_bytes = output_buffer.getvalue()
            content_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif file_ext == 'xlsx':
            for mapping in field_mappings:
                mapping['field']['obj'].value = mapping['value']
            output_buffer = io.BytesIO()
            wb.save(output_buffer)
            file_bytes = output_buffer.getvalue()
            content_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif file_ext == 'pptx':
            for mapping in field_mappings:
                shape = mapping['field']['obj']
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.text = mapping['value']
            output_buffer = io.BytesIO()
            prs.save(output_buffer)
            file_bytes = output_buffer.getvalue()
            content_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        output_filename = target_filename.rsplit('.', 1)[0] + '_filled.' + file_ext
        filled_key = f"document-autofill/completed/{session_id}_{output_filename}"
        s3_client.put_object(Bucket=BUCKET, Key=filled_key, Body=file_bytes, ContentType=content_type)
        download_url = s3_client.generate_presigned_url('get_object', Params={'Bucket': BUCKET, 'Key': filled_key}, ExpiresIn=3600)
        logger.info(f"Document filled: {filled_key} ({len(field_mappings)} fields)")
        return cors_response({"downloadUrl": download_url, "filename": output_filename})
    except Exception as e:
        logger.error(f"Fill failed: {e}", exc_info=True)
        return cors_response({"error": str(e)}, 500)








