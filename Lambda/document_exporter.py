"""
Document Export Module
Generates Excel, Word, and PowerPoint documents from agent responses
"""

import io
import re
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor


def detect_content_type(text):
    """Detect if text contains tables, lists, or is plain text."""
    # Check for markdown tables
    has_table = bool(re.search(r'\|.*\|.*\|', text))
    # Check for bullet points or numbered lists
    has_list = bool(re.search(r'(^|\n)[\*\-\d]+[\.\)]\s', text))
    return {
        'has_table': has_table,
        'has_list': has_list,
        'is_plain_text': not has_table and not has_list
    }


def parse_markdown_table(text):
    """Parse markdown table into rows and columns."""
    lines = text.strip().split('\n')
    table_data = []
    
    for line in lines:
        if '|' in line:
            # Skip separator lines (e.g., |---|---|)
            if re.match(r'^\s*\|[\s\-:]+\|\s*$', line):
                continue
            # Extract cells
            cells = [cell.strip() for cell in line.split('|')[1:-1]]
            if cells:
                table_data.append(cells)
    
    return table_data


def export_to_excel(text, title="Export"):
    """Generate Excel file from text (tables or lists)."""
    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]  # Excel sheet name limit
    
    content_type = detect_content_type(text)
    
    if content_type['has_table']:
        # Parse and export table
        table_data = parse_markdown_table(text)
        
        if table_data:
            # Add header row with styling
            for col_idx, cell_value in enumerate(table_data[0], 1):
                cell = ws.cell(row=1, column=col_idx, value=cell_value)
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                cell.alignment = Alignment(horizontal="center", vertical="center")
            
            # Add data rows
            for row_idx, row_data in enumerate(table_data[1:], 2):
                for col_idx, cell_value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                    cell.alignment = Alignment(horizontal="left", vertical="center")
            
            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                ws.column_dimensions[column_letter].width = min(max_length + 2, 50)
    else:
        # Export as simple list
        lines = text.strip().split('\n')
        for row_idx, line in enumerate(lines, 1):
            ws.cell(row=row_idx, column=1, value=line.strip())
        ws.column_dimensions['A'].width = 80
    
    # Save to BytesIO
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output


def export_to_word(text, title="Export"):
    """Generate Word document from text."""
    doc = Document()
    
    # Add title
    title_para = doc.add_heading(title, level=1)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Parse content
    content_type = detect_content_type(text)
    
    if content_type['has_table']:
        # Add table
        table_data = parse_markdown_table(text)
        if table_data:
            table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
            table.style = 'Light Grid Accent 1'
            
            # Fill table
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.rows[row_idx].cells[col_idx]
                    cell.text = cell_value
                    # Bold header row
                    if row_idx == 0:
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True
    else:
        # Add paragraphs
        lines = text.strip().split('\n')
        for line in lines:
            if line.strip():
                # Detect headers (lines ending with :)
                if line.strip().endswith(':'):
                    doc.add_heading(line.strip(), level=2)
                else:
                    para = doc.add_paragraph(line.strip())
                    para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    
    # Save to BytesIO
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def export_to_powerpoint(text, title="Presentation"):
    """Generate PowerPoint presentation from text."""
    prs = Presentation()
    
    # Set slide size (16:9)
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = "Generated by AI Document Assistant"
    
    # Parse content into sections
    lines = text.strip().split('\n')
    current_section = None
    current_bullets = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Detect section headers (lines ending with : or starting with ##)
        if line.endswith(':') or line.startswith('##'):
            # Save previous section
            if current_section and current_bullets:
                add_content_slide(prs, current_section, current_bullets)
            # Start new section
            current_section = line.rstrip(':').lstrip('#').strip()
            current_bullets = []
        else:
            # Add to current section
            current_bullets.append(line)
    
    # Add last section
    if current_section and current_bullets:
        add_content_slide(prs, current_section, current_bullets)
    
    # If no sections detected, create single content slide
    if len(prs.slides) == 1:
        add_content_slide(prs, "Content", lines)
    
    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def add_content_slide(prs, title, bullets):
    """Add a content slide with title and bullet points."""
    bullet_slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Add bullets (max 6 per slide for readability)
    for bullet in bullets[:6]:
        p = text_frame.add_paragraph()
        p.text = bullet.lstrip('*-•').strip()
        p.level = 0
        p.font.size = PptxPt(18)
